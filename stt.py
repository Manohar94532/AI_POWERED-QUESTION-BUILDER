from streamlit_lottie import st_lottie
from streamlit_option_menu import option_menu
import requests
import streamlit as st
import streamlit_chat as stc
import google.generativeai as genai
import PyPDF2
import docx
import os
import platform
import random
# Changed mysql.connector to pymongo
from pymongo import MongoClient
from pymongo.errors import ConnectionFailure, OperationFailure
from bson.objectid import ObjectId  # Import ObjectId for MongoDB's _id
from werkzeug.security import generate_password_hash, check_password_hash
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import nltk
from nltk.sentiment import SentimentIntensityAnalyzer
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import re
import json
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D
import seaborn as sns
import googletrans
import numpy as np
import io
import datetime
from datetime import datetime, timedelta
from fpdf import FPDF
import csv
import tempfile
from pptx import Presentation
from streamlit_option_menu import option_menu
from pptx.util import Inches
import plotly.express as px
from streamlit_lottie import st_lottie  # Import the Lottie function
import requests  # To fetch the Lottie animation
import googletrans
#from google_trans_new import google_translator
# Using deep-translator as it's more reliable
from deep_translator import GoogleTranslator


# Initialize translator (from google_trans_new, if used elsewhere)
#translator = google_translator()


st.set_page_config(layout="wide")


# --- TRAINER DASHBOARD ---
def trainer_dashboard():
    trainer_username = st.session_state.user['username']

    # Session state initialization
     # --- FIX: INITIALIZE CHAT HISTORY HERE ---
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    # -------------------------------------------
    
    # Session state initialization for questions
    if 'generated_questions' not in st.session_state:
        st.session_state.generated_questions = []

    with st.sidebar:
        selected = option_menu(
            menu_title="Trainer Dashboard",
            options=[
                "Assign Employees", "Upload Curriculum", "Generate Question Bank", "View Questions",
                "Assign Learning Plan", "Review Feedback", "Employee Performance", "Chatbot"
            ],
            icons=[
                "person-plus", "upload", "question-circle", "eye", "card-checklist",
                "check2-circle", "bar-chart-line", "chat-dots"
            ],
            menu_icon="briefcase", default_index=0
        )

     # Initialize question_banks to avoid UnboundLocalError
    question_banks = []
    # Display content based on the selected option

    # --- NEW: Assign Employees Page ---
    if selected == "Assign Employees":
        st.subheader("Assign New Employees 🧑‍🏫")
        unassigned_employees = get_unassigned_employees()

        if not unassigned_employees:
            st.info("There are currently no unassigned employees.")
        else:
            employee_usernames = [emp['username']
                                  for emp in unassigned_employees]
            selected_employees = st.multiselect(
                "Select employees to assign to yourself:",
                options=employee_usernames
            )
            if st.button("Assign Selected Employees"):
                if not selected_employees:
                    st.warning("Please select at least one employee.")
                elif assign_employees_to_trainer(selected_employees, trainer_username):
                    st.success(
                        f"Successfully assigned {', '.join(selected_employees)}!")
                    st.rerun()
                else:
                    st.error("An error occurred during assignment.")

    elif selected == "Upload Curriculum":
        st.subheader("Upload Curriculum 📁")
        technology = st.text_input("Technology")
        topics = st.text_area("Topics (one per line)")
        uploaded_file = st.file_uploader("Upload curriculum file", type=[
                                         'pdf', 'docx', 'txt', 'pptx'])

        if st.button("Upload Curriculum"):
            if technology and uploaded_file:
                content = extract_text_from_file(uploaded_file)
                topic_list = [t.strip()
                              for t in topics.split('\n') if t.strip()]
                if upload_curriculum(technology, topic_list, content, trainer_username):
                    st.success("Curriculum uploaded successfully!")
                else:
                    st.error("Failed to upload curriculum.")
            else:
                st.warning(
                    "Please provide a technology name and upload a file.")

    # --- Generate Question Bank Page (Modified) ---
    elif selected == "Generate Question Bank":
        st.subheader("Generate Question Bank 📚")
        curricula = get_trainer_curricula(trainer_username)
        if not curricula:
            st.warning(
                "No curricula available. Please upload a curriculum first.")
        else:
            curriculum_map = {c['technology']: c['content'] for c in curricula}
            selected_tech = st.selectbox(
                "Select Curriculum", options=list(curriculum_map.keys()))

            num_questions = st.number_input(
                "Number of Questions", min_value=1, value=5)
            question_type = st.selectbox(
                "Question Type", ["multiple-choice", "subjective", "fill-in-the-blank"])
            difficulty = st.selectbox("Difficulty", ["Easy", "Medium", "Hard"])

            if st.button("Generate Question Bank"):
                content = curriculum_map.get(selected_tech)
                if content:
                    questions, options, correct_answers = generate_questions(
                        content, num_questions, question_type)
                    question_bank_id = save_question_bank(
                        selected_tech, [], '\n'.join(questions), difficulty,
                        '\n'.join(correct_answers), question_type,
                        '|||'.join(['###'.join(opt) for opt in options]),
                        trainer_username
                    )
                    if question_bank_id:
                        st.success(
                            f"Question Bank generated successfully! ID: {question_bank_id}")
                    else:
                        st.error("Failed to save question bank.")
                else:
                    st.error("Failed to retrieve curriculum content.")

    # --- View Questions Page (Modified) ---
    elif selected == "View Questions":
        st.subheader("View Your Question Banks 📖")
        question_banks = get_trainer_question_banks(trainer_username)
        if not question_banks:
            st.info("You have not created any question banks yet.")
        else:
            qb_options = {str(
                qb['_id']): f"{qb['technology']} - {qb['difficulty']}" for qb in question_banks}
            selected_id = st.selectbox("Select Question Bank", options=qb_options.keys(
            ), format_func=lambda x: qb_options[x])

            if selected_id:
                qb_details = next(
                    (qb for qb in question_banks if str(qb['_id']) == selected_id), None)
                if qb_details:
                    st.info(
                        f"**Technology:** {qb_details['technology']} | **Difficulty:** {qb_details['difficulty']}")
                    st.write("---")
                    questions = qb_details.get('questions', '').split('\n')
                    for i, q in enumerate(questions, 1):
                        st.write(f"**Q{i}:** {q}")

    # --- Assign Learning Plan Page (Modified) ---
    # In trainer_dashboard(), find the "Assign Learning Plan" block and update the button logic:

    elif selected == "Assign Learning Plan":
        st.subheader("Assign Learning Plan to Your Employees 🧑‍🏫")
        assigned_employees = get_assigned_employees(trainer_username)
        question_banks = get_trainer_question_banks(trainer_username)

        if not assigned_employees:
            st.warning("You have no employees assigned. Please assign them from the 'Assign Employees' tab.")
        elif not question_banks:
            st.warning("You have no question banks. Please create one first.")
        else:
            employee_options = [emp['username'] for emp in assigned_employees]
            selected_employee = st.selectbox("Select Employee", options=employee_options)

            qb_options = {str(qb['_id']): f"{qb['technology']} - {qb['difficulty']}" for qb in question_banks}
            selected_qb_id = st.selectbox("Select Question Bank to Assign", options=list(qb_options.keys()), format_func=lambda x: qb_options[x])

            if st.button("Assign Plan"):
                # --- THIS IS THE CORRECTED FUNCTION CALL ---
                plan_id = create_learning_plan(ObjectId(selected_qb_id), selected_employee, trainer_username)
                if plan_id:
                    st.success(f"Learning plan '{qb_options[selected_qb_id]}' assigned to {selected_employee}!")
                    send_notification(
                        recipient_role="employee",
                        message=f"New plan '{qb_options[selected_qb_id]}' was assigned by {trainer_username}.",
                        username=selected_employee
                    )
                else:
                    st.error("Failed to assign plan. It may already be assigned.")

    elif selected == "Review Feedback":
        st.subheader("Review Feedback 🔍")
        feedback = review_feedback()
        if not feedback:
            st.info("No feedback available yet.")
        else:
            feedback_df = pd.DataFrame(feedback)
            st.dataframe(feedback_df)

            # Debugging: Check the structure of question_banks
            question_banks = get_all_question_banks()  # Ensure question_banks is populated
            st.write("Question Banks DataFrame:")
            # This will show the structure of the DataFrame
            st.write(pd.DataFrame(question_banks))

            # Check if 'id' exists in question_banks (or _id for MongoDB)
            if not question_banks or '_id' not in question_banks[0]:
                st.error("The '_id' column is missing from question_banks.")
                return

            # Sentiment Analysis Summary
            sentiment_counts = feedback_df['sentiment'].value_counts()
            st.subheader("Sentiment Analysis Summary")
            st.bar_chart(sentiment_counts)

            # Question Bank Feedback Summary
            st.subheader("Question Bank Feedback Summary")
            # Convert ObjectId to string for merging
            feedback_df['question_bank_id_str'] = feedback_df['question_bank_id'].apply(
                str)
            qb_df = pd.DataFrame(question_banks)
            qb_df['id_str'] = qb_df['_id'].apply(str)

            qb_feedback = feedback_df.groupby('question_bank_id_str')[
                'rating'].mean().reset_index()
            qb_feedback = qb_feedback.merge(
                qb_df, left_on='question_bank_id_str', right_on='id_str')
            st.dataframe(
                qb_feedback[['question_bank_id_str', 'technology', 'difficulty', 'rating']])

    elif selected == "Curriculum Overview":
        st.subheader("Curriculum Overview 📜")
        curricula = get_all_curricula()

        if curricula:
            # Convert ASCII values to strings for each curriculum (if applicable)
            for curriculum in curricula:
                # MongoDB stores topics as strings or lists directly, so no ASCII conversion needed unless explicitly stored as ASCII values
                if isinstance(curriculum.get('topics'), list) and all(isinstance(x, int) for x in curriculum['topics']):
                    curriculum['topics'] = ascii_to_string(
                        curriculum['topics'])  # Convert to string

            curriculum_df = pd.DataFrame(curricula)
            st.dataframe(curriculum_df[['technology']])
        else:
            st.info("No curricula available.")

        # Display feedback summary for question banks
        st.subheader("Question Bank Feedback Summary")
        feedback = review_feedback()  # Assuming this function retrieves feedback data

        if feedback:
            feedback_df = pd.DataFrame(feedback)
            st.dataframe(feedback_df)  # Display feedback DataFrame
        else:
            st.info("No feedback available yet.")

        st.subheader("Clear History")
        if st.button("Clear Curriculum Content History"):
            db = create_connection()
            if db is not None:
                try:
                    # Clear the curriculum collection
                    db.curriculum.delete_many({})
                    st.success(
                        "Curriculum content history cleared successfully!")
                except OperationFailure as e:
                    st.error(f"Failed to clear curriculum history: {e}")
            else:
                st.error("Failed to connect to database")

        if st.button("Clear Generated Topics History"):
            db = create_connection()
            if db is not None:
                try:
                    # Clear the generated_question_files collection
                    db.generated_question_files.delete_many({})
                    st.success(
                        "Generated topics history cleared successfully!")
                except OperationFailure as e:
                    st.error(f"Failed to clear generated topics history: {e}")
            else:
                st.error("Failed to connect to database")

    # In trainer_dashboard()
    elif selected == "Download Questions":
        st.subheader("Download Questions ⬇️")
        question_banks = get_all_question_banks()
        if not question_banks:
            st.info("No question banks available yet.")
        else:
            selected_qb = st.selectbox("Select Question Bank", options=[(str(
                qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks], format_func=lambda x: f"ID: {x[0]} - {x[1]}", key="download_qb_select")

            # CORRECTED: All logic is now nested inside this 'if' block.
            if selected_qb:
                qb_id_str, _ = selected_qb
                qb_id = ObjectId(qb_id_str)
                qb_details = next(
                    (qb for qb in question_banks if qb['_id'] == qb_id), None)

                if qb_details:
                    st.write(f"Technology: {qb_details['technology']}")
                    st.write(f"Difficulty: {qb_details['difficulty']}")

                    file_format = st.selectbox("Select File Format", [
                                               "docx", "pdf", "pptx", "csv"])
                    questions = qb_details['questions'].split('\n')

                    if file_format == "pdf":
                        pdf = FPDF()
                        pdf.add_page()
                        pdf.set_font("Arial", size=12)
                        for q in questions:
                            # Encode to prevent errors with special characters
                            pdf.multi_cell(0, 10, txt=q.encode(
                                'latin-1', 'replace').decode('latin-1'))

                        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp:
                            pdf.output(temp.name)
                            temp.seek(0)
                            st.download_button(
                                label='Download PDF',
                                data=temp.read(),
                                file_name=f'questions_{qb_id_str}.pdf',
                                mime='application/pdf'
                            )

                    elif file_format == "docx":
                        doc = docx.Document()
                        for question in questions:
                            doc.add_paragraph(question)
                        buffer = io.BytesIO()
                        doc.save(buffer)
                        buffer.seek(0)
                        st.download_button(
                            label='Download DOCX',
                            data=buffer.getvalue(),
                            file_name=f'questions_{qb_id_str}.docx',
                            mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                        )

                    elif file_format == "pptx":
                        presentation = Presentation()
                        slide_layout = presentation.slide_layouts[6]
                        slide = presentation.slides.add_slide(slide_layout)
                        left, top, height = Inches(1), Inches(1), Inches(0.5)

                        for question in questions:
                            textbox = slide.shapes.add_textbox(
                                left, top, width=Inches(8), height=height)
                            textbox.text = question
                            top += height + Inches(0.1)

                        buffer = io.BytesIO()
                        presentation.save(buffer)
                        buffer.seek(0)
                        st.download_button(
                            label='Download PPTX',
                            data=buffer.getvalue(),
                            file_name=f'questions_{qb_id_str}.pptx',
                            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        )

                    elif file_format == "csv":
                        buffer = io.StringIO()
                        csv_writer = csv.writer(buffer)
                        csv_writer.writerow(["Question"])
                        for question in questions:
                            csv_writer.writerow([question])
                        buffer.seek(0)
                        st.download_button(
                            label='Download CSV',
                            data=buffer.getvalue(),
                            file_name=f'questions_{qb_id_str}.csv',
                            mime='text/csv'
                        )

    elif selected == "Employee Performance":
        st.subheader("Employee Performance 📈")
        employees = get_all_users()

        if employees:
            selected_employee = st.selectbox(
                "Select Employee",
                options=[employee['username'] for employee in employees],
                key="employee_performance_select"
            )

            if selected_employee:
                # Fetch assessment results for the selected employee
                assessment_results = get_assessment_results(selected_employee)
                if assessment_results:
                    # Prepare data for the table
                    performance_data = []
                    for result in assessment_results:
                        performance_data.append({
                            # Convert ObjectId to string
                            'Question Bank ID': str(result['question_bank_id']),
                            'Score': result['score'],
                            'Completed At': result['completed_at']
                        })

                    # Convert to DataFrame for better visualization
                    performance_df = pd.DataFrame(performance_data)

                    # Display summary metrics
                    st.subheader(f"Summary Statistics for {selected_employee}")
                    total_assessments = len(performance_df)
                    avg_score = performance_df['Score'].mean(
                    ) if total_assessments > 0 else 0
                    best_score = performance_df['Score'].max(
                    ) if total_assessments > 0 else 0

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Assessments", total_assessments)
                    with col2:
                        st.metric("Average Score", f"{avg_score:.1f}")
                    with col3:
                        st.metric("Best Score", best_score)

                    # Display the performance data in a styled table
                    st.write(f"Performance Data for {selected_employee}:")
                    st.dataframe(performance_df.style.highlight_max(
                        axis=0))  # Highlight max scores

                    # Convert 'Completed At' column to datetime for sorting
                    performance_df['Completed At'] = pd.to_datetime(
                        performance_df['Completed At'])
                    performance_df.sort_values('Completed At', inplace=True)

                    # Create visualizations
                    st.subheader("Performance Over Time")

                    # Line chart for scores over time
                    fig_line = px.line(performance_df, x='Completed At', y='Score',
                                       title='Score Over Time', markers=True)
                    st.plotly_chart(fig_line)

                    # Bar chart for scores by question bank
                    fig_bar = px.bar(performance_df, x='Question Bank ID', y='Score',
                                     title='Scores by Question Bank', text='Score')
                    st.plotly_chart(fig_bar)

                    # Convert figures to HTML format for download
                    fig_line_html = fig_line.to_html(full_html=False)
                    fig_bar_html = fig_bar.to_html(full_html=False)

                    # Provide download buttons
                    st.download_button(label="Download Line Chart as HTML", data=fig_line_html,
                                       file_name=f"{selected_employee}_performance_over_time.html", mime="text/html")

                    st.download_button(label="Download Bar Chart as HTML", data=fig_bar_html,
                                       file_name=f"{selected_employee}_score_by_question_bank.html", mime="text/html")

                    st.download_button(label="Download Performance Data as CSV",
                                       data=performance_df.to_csv(index=False),
                                       file_name=f"{selected_employee}_performance.csv", mime="text/csv")

                else:
                    st.info("No assessment results available for this employee.")
        else:
            st.info("No employees available.")

    # Display content based on the selected option
    if selected == "Generate Questions":
        # Horizontal menu for question generation methods
        question_generation_method = option_menu(
            menu_title=None,  # required
            options=["Generate Questions by Topic",
                     "Generate Questions from Prompt"],  # required
            icons=["book", "pencil"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="horizontal",
        )

        if question_generation_method == "Generate Questions by Topic":
            st.subheader("Generate Questions by Topic 🚀")
            topic_name = st.text_input(
                "Enter Topic Name", key="topic_input")  # Input for topic name
            num_questions = st.number_input("Number of Questions to Generate", min_value=1,
                                            value=5, key="num_questions_input")  # Input for number of questions

            # Dropdown for selecting question type
            question_type = st.selectbox("Select Question Type", [
                                         "Multiple Choice", "Subjective", "Fill in the Blanks"], key="question_type_select")

            if st.button("Generate Questions", key="generate_topic_questions_button"):
                if topic_name:
                    try:
                        prompt = f"Generate {num_questions} {question_type.lower()} questions based on the topic: {topic_name}."
                        generated_questions = model.generate_content(
                            prompt)  # Using the model to generate questions
                        questions_text = generated_questions.text.strip()  # Extracting the generated text

                        # Displaying the generated questions
                        st.write("Generated Questions:")
                        questions_list = questions_text.split('\n')
                        selected_questions = []

                        for i, question in enumerate(questions_list, 1):
                            if question.strip():  # Only show non-empty questions
                                # Display each question with a checkbox
                                if st.checkbox(f"Question {i}: {question.strip()}", value=True):
                                    selected_questions.append(question.strip())

                        # Store generated questions in session state
                        if 'history' not in st.session_state:
                            st.session_state.history = []
                        st.session_state.history.append(
                            {"topic": topic_name, "questions": selected_questions})

                        st.session_state.generated_questions = selected_questions
                        st.session_state.topic_name = topic_name
                        st.success(
                            f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                    except Exception as e:
                        st.error(f"Error generating questions: {e}")
                else:
                    st.error("Please enter a topic name.")

            # New feature: Add Questions to Question Bank
            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")

                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [
                    f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox(
                    "Select Question Bank", options=qb_options)

                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input(
                            "Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", [
                                                         "Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(
                                new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(
                                    f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split(
                            # Convert to ObjectId
                            '-')[0].split(':')[1].strip())
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(
                                f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error(
                                "Failed to add questions to the selected question bank.")

                    # Clear the generated questions from session state
                    del st.session_state.generated_questions
                    del st.session_state.topic_name

        elif question_generation_method == "Generate Questions from Prompt":
            st.subheader("Generate Questions from Prompt ✍️")
            topic_name = st.text_input("Enter Topic Name")
            prompt = st.text_area("Enter a paragraph to generate questions")
            question_type = st.selectbox("Select Question Type", [
                                         "Multiple Choice", "Subjective", "Fill in the Blanks"])
            difficulty_level = st.selectbox("Select Difficulty Level", [
                                            "Easy", "Medium", "Hard"])
            num_questions = st.number_input(
                "Number of Questions to Generate", min_value=1, value=10)

            if st.button("Generate Questions"):
                generated_questions = generate_questions_from_prompt(
                    prompt, question_type, difficulty_level, num_questions, topic_name)
                if generated_questions:
                    st.write("Generated Questions:")
                    selected_questions = []
                    for i, question in enumerate(generated_questions):
                        if st.checkbox(f"Question {i+1}", value=True):
                            selected_questions.append(question)
                        st.write(question)

                    # Store generated questions in session state
                    if 'history' not in st.session_state:
                        st.session_state.history = []
                    st.session_state.history.append(
                        {"topic": topic_name, "questions": selected_questions})

                    st.session_state.generated_questions = selected_questions
                    st.session_state.topic_name = topic_name
                    st.success(
                        f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                else:
                    st.error("Failed to generate questions")

            # New feature: View Generated Questions History
            st.subheader("View Generated Questions History")
            topic_questions = get_generated_questions_history()  # Retrieve the history

            if topic_questions:
                selected_topic = st.selectbox(
                    "Select Topic", options=list(topic_questions.keys()))
                if selected_topic:
                    questions = topic_questions[selected_topic]
                    st.write(
                        f"Generated Questions for Topic: {selected_topic}")
                    for i, question in enumerate(questions, 1):
                        st.write(f"{i}. {question}")
            else:
                st.info("No generated questions history available.")

            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")

                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [
                    f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox(
                    "Select Question Bank", options=qb_options)

                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input(
                            "Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", [
                                                         "Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(
                                new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(
                                    f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split(
                            # Convert to ObjectId
                            '-')[0].split(':')[1].strip())
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(
                                f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error(
                                "Failed to add questions to the selected question bank.")

        # Clear the generated questions from session state
                        del st.session_state.generated_questions
                        del st.session_state.topic_name
        # Display content based on the selected option

    elif selected == "Chatbot":
        # Display chatbot interface at the top
        st.subheader("Chat with the AI Trainer 🤖")

        # Initialize the message state if not exists
        if "msg" not in st.session_state:
            st.session_state.msg = ""

        # Create a container for the chat interface
        chat_container = st.container()

        # Define avatars
        user_avatar = "https://static.vecteezy.com/system/resources/previews/009/664/418/non_2x/people-user-team-transparent-free-png.png"
        ai_avatar = "https://thumbs.dreamstime.com/b/chatbot-logo-messenger-ai-robot-icon-vector-illustration-277900892.jpg"

        def clear_text():
            st.session_state.msg = st.session_state.user_input
            st.session_state.user_input = ""

        with chat_container:
            # Display chat messages
            for chat in st.session_state.chat_history:
                if chat['role'] == 'assistant':
                    # Chatbot message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;color:black'>"
                        f"<img src='{ai_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-right: 10px;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>AI:</strong> {chat['content']}</div></div>",
                        unsafe_allow_html=True
                    )
                else:
                    # User message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0; justify-content: flex-end;color:black'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {chat['content']}</div>"
                        f"<img src='{user_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-left: 10px;'>"
                        f"</div>",
                        unsafe_allow_html=True
                    )

            # Input field for user to enter a prompt
            st.text_input("Type your message here...", key="user_input",
                          placeholder="Type a message...", on_change=clear_text)

            if st.session_state.msg:  # Only process if there's a message
                # Append user input to chat history
                st.session_state.chat_history.append(
                    {"role": "user", "content": st.session_state.msg})

                try:
                    # Generate AI response
                    prompt = f"You are an AI assistant for trainers. Respond to the following message: {st.session_state.msg}"
                    response = model.generate_content(prompt)

                    # Handle the response properly for Gemini model
                    if hasattr(response, 'parts'):
                        ai_response = ''.join(
                            part.text for part in response.parts)
                    else:
                        ai_response = response.candidates[0].content.parts[0].text

                    # Append AI response to chat history
                    st.session_state.chat_history.append(
                        {"role": "assistant", "content": ai_response})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    ai_response = "I apologize, but I encountered an error. Please try again."
                    st.session_state.chat_history.append(
                        {"role": "assistant", "content": ai_response})

                # Clear the message state
                st.session_state.msg = ""

                # Rerun the app to display the new messages
                st.rerun()

    # Get notifications for trainer
    notifications = get_notifications("trainer", None)
    # Display notifications in the sidebar
    display_notifications(notifications, "trainer")

    if notifications:
        st.sidebar.write("Notifications:")
        for notification in notifications:
            st.sidebar.write(notification['message'])
    else:
        st.sidebar.write("No notifications available.")


# Download NLTK data
try:
    nltk.data.find('sentiment/vader_lexicon.zip')
except LookupError:
    st.info("Downloading NLTK 'vader_lexicon' for sentiment analysis...")
    nltk.download('vader_lexicon')
    st.success("Download complete.")
# Load .env file (if used locally for development)
# from dotenv import load_dotenv
# load_dotenv()

# Set Google API Key from environment variable
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    st.error("Google API key not found. Please set the GOOGLE_API_KEY environment variable.")
else:
    genai.configure(api_key=GOOGLE_API_KEY)

# Configure Google Generative AI model
genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
model = genai.GenerativeModel('gemini-2.0-flash')


# Modified save_question_bank to use MongoDB
def save_question_bank(technology, topics, questions, difficulty, correct_answers, question_type, options, trainer_username):
    """Saves a question bank scoped to a specific trainer."""
    db = create_connection()
    if db is None:
        return None
    try:
        qb_doc = {
            "trainer_username": trainer_username, # <-- Saves the trainer's name
            "technology": technology,
            "topics": topics,
            "questions": questions,
            "difficulty": difficulty,
            "question_type": question_type,
            "options": options,
            "created_at": datetime.now()
        }
        result_qb = db.question_banks.insert_one(qb_doc)
        question_bank_id = result_qb.inserted_id

        answer_doc = {
            "question_bank_id": question_bank_id,
            "answer_data": correct_answers
        }
        db.question_answers.insert_one(answer_doc)
        return str(question_bank_id)
    except OperationFailure as e:
        st.error(f"MongoDB operation error: {e}")
        return None

# MongoDB connection


@st.cache_resource
def create_connection():
    """Create and cache a MongoDB database connection."""
    mongo_uri = os.getenv("MONGO_URI")
    if not mongo_uri:
        st.error(
            "MongoDB URI not found. Please set the MONGO_URI environment variable.")
        return None
    try:
        client = MongoClient(mongo_uri)
        db = client["final_mongodb"]
        return db
    except ConnectionFailure as e:
        st.error(f"Error connecting to MongoDB: {e}")
        return None


db = create_connection()


def get_unassigned_employees():
    """Fetches employees who are not yet assigned to any trainer."""
    if db is None:
        return []
    try:
        unassigned = db.users.find({
            "role": "Employee",
            "assigned_trainer": None
        })
        return list(unassigned)
    except OperationFailure as e:
        st.error(f"DB Error fetching unassigned employees: {e}")
        return []


def assign_employees_to_trainer(employee_usernames, trainer_username):
    """Assigns a list of employees to the specified trainer."""
    if db is None:
        return False
    try:
        result = db.users.update_many(
            {"username": {"$in": employee_usernames}},
            {"$set": {"assigned_trainer": trainer_username}}
        )
        return result.modified_count > 0
    except OperationFailure as e:
        st.error(f"DB Error assigning employees: {e}")
        return False


def get_assigned_employees(trainer_username):
    """Fetches all employees assigned to a specific trainer."""
    if db is None:
        return []
    try:
        employees = db.users.find({
            "role": "Employee",
            "assigned_trainer": trainer_username
        })
        return list(employees)
    except OperationFailure as e:
        st.error(f"DB Error fetching assigned employees: {e}")
        return []


def ascii_to_string(ascii_list):
    """Convert a list of ASCII values to a string."""
    return ''.join(chr(num) for num in ascii_list)

    # Example usage when retrieving topics
    curricula = get_all_curricula()
    for curriculum in curricula:
        # Assuming 'topics' is a list of ASCII values in the curricula
        # Check if topics is a list of ASCII values
        if isinstance(curriculum['topics'], list):
            curriculum['topics'] = ascii_to_string(
                curriculum['topics'])  # Convert to string


def format_timestamp(timestamp_str):
    """
    Format timestamp for display
    """
    if timestamp_str == 'N/A':
        return 'N/A'
    try:
        # Parse the timestamp string to datetime
        timestamp = datetime.datetime.strptime(
            timestamp_str, "%Y-%m-%d %H:%M:%S")
        # Format it to a more readable format
        return timestamp.strftime("%b %d, %Y %I:%M %p")
    except (ValueError, TypeError):
        return 'N/A'


# Utility functions
def extract_text_from_file(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    text = ""

    try:
        if file_extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_extension == '.docx':
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == '.txt':
            text = file.getvalue().decode('utf-8')
        elif file_extension in ['.ppt', '.pptx']:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        elif file_extension == '.csv':
            csv_data = pd.read_csv(file)
            text = csv_data.to_string(index=False)
        else:
            # For other file types, attempt to read as text
            try:
                text = file.getvalue().decode('utf-8')
            except UnicodeDecodeError:
                raise ValueError(
                    f"Unable to extract text from {file_extension} file.")
    except Exception as e:
        raise ValueError(f"Error processing {file_extension} file: {str(e)}")

    # Clean the extracted text
    cleaned_text = clean_text(text)
    return cleaned_text


def clean_text(text):
    # Remove non-printable characters and control characters
    # Keep only printable ASCII characters
    text = re.sub(r'[^\x20-\x7E]+', ' ', text)

    # Normalize whitespace
    # Replace multiple spaces with a single space
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()  # Trim leading and trailing whitespace

    return text


def create_new_question_bank(technology, difficulty, questions):
    db = create_connection()
    if db is None:
        return None

    try:
        questions_text = '\n'.join(questions)
        new_qb_doc = {
            "technology": technology,
            "difficulty": difficulty,
            "questions": questions_text,
            "created_at": datetime.now()
        }
        result = db.question_banks.insert_one(new_qb_doc)
        return str(result.inserted_id)  # Return the ObjectId as a string
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


import random
import time

def generate_questions(text, num_questions=5, question_type="multiple-choice"):
    # Add randomization to ensure unique questions each time
    random_seed = random.randint(1000, 9999)
    timestamp = int(time.time()) % 10000
    
    if question_type == "multiple-choice":
        # --- ENHANCED PROMPT WITH UNIQUENESS FACTORS ---
        prompt = f"""Generate exactly {num_questions} unique multiple-choice questions based on the following text. 

        Important: Create completely new and different questions each time. Focus on different aspects, concepts, and details from the text.

        Text to analyze:
        {text}

        Requirements:
        - Generate exactly {num_questions} questions
        - Each question should test different aspects of the content
        - Provide 4 options (A, B, C, D) for each question
        - Only ONE option should be correct
        - Include the correct answer letter after each question

        Format (strictly follow this format):
        
        Q1: [Question about specific concept/detail]
        A) [Option 1]
        B) [Option 2] 
        C) [Option 3]
        D) [Option 4]
        Answer: [Correct letter]

        Q2: [Question about different concept/detail]
        A) [Option 1]
        B) [Option 2]
        C) [Option 3] 
        D) [Option 4]
        Answer: [Correct letter]

        Generation ID: {random_seed}-{timestamp}
        """
    elif question_type == "subjective":
        prompt = f"""Generate exactly {num_questions} unique subjective questions based on the following text. 
        
        Create diverse questions that explore different aspects of the content:
        
        {text}
        
        Generation ID: {random_seed}-{timestamp}
        """
    elif question_type == "fill-in-the-blank":
        prompt = f"""Generate exactly {num_questions} unique fill-in-the-blank questions based on the following text. 
        
        Create questions that test different key terms and concepts:
        
        {text}
        
        Provide each answer on a new line starting with 'A:'.
        
        Generation ID: {random_seed}-{timestamp}
        """
    else:
        raise ValueError("Invalid question type")

    # Configure model for more randomness (if using Google Gemini)
    generation_config = {
        'temperature': 0.9,  # Higher temperature for more randomness
        'top_p': 0.8,
        'top_k': 40,
        'max_output_tokens': 2048,
    }

    try:
        # Generate with randomness config
        response = model.generate_content(
            prompt,
            generation_config=generation_config
        )
        generated_text = response.text
    except:
        # Fallback if generation_config is not supported
        response = model.generate_content(prompt)
        generated_text = response.text

    print(f"Generating {num_questions} {question_type} questions... ID: {random_seed}-{timestamp}")
    
    questions = []
    options = []
    correct_answers = []

    # Split the entire text by the question marker "Q" followed by a number and colon.
    question_blocks = re.split(r'Q\d+:', generated_text)[1:]

    for i, block in enumerate(question_blocks):
        if not block.strip() or i >= num_questions:  # Ensure we don't exceed requested number
            continue

        lines = [line.strip() for line in block.strip().split('\n') if line.strip()]

        if not lines:
            continue

        # --- ROBUST PARSING LOGIC ---
        question_text = lines[0]
        # Clean up question text
        question_text = question_text.replace('Generation ID:', '').strip()
        if question_text:
            questions.append(question_text)

            if question_type == "multiple-choice":
                current_options = []
                answer_line = ""
                
                for line in lines[1:]:
                    if re.match(r'^[A-D]\)', line):  # Matches A), B), C), D)
                        option_text = line.split(') ', 1)[1] if ') ' in line else line[2:].strip()
                        current_options.append(option_text)
                    elif line.lower().startswith('answer:'):
                        answer_line = line
                    elif line.startswith('Generation ID:'):
                        break  # Stop parsing when we hit the generation ID

                options.append(current_options)

                # --- Store the correct answer LETTER ---
                if answer_line and current_options:
                    try:
                        correct_letter = answer_line.split(':')[1].strip().upper()
                        # Validate that it's a valid option letter
                        if correct_letter in ['A', 'B', 'C', 'D'] and ord(correct_letter) - ord('A') < len(current_options):
                            correct_answers.append(correct_letter)
                        else:
                            print(f"Warning: Invalid answer letter '{correct_letter}' for question: {question_text}")
                            correct_answers.append('A')  # Default fallback
                    except Exception as e:
                        print(f"Error parsing answer for question '{question_text}': {e}")
                        correct_answers.append('A')  # Default fallback
                else:
                    print(f"Warning: No answer found for question: {question_text}")
                    correct_answers.append('A')  # Default fallback

            elif question_type == "fill-in-the-blank":
                options.append([])  # No options for fill-in-the-blank
                if len(lines) > 1 and lines[1].startswith('A:'):
                    correct_answers.append(lines[1].split('A:', 1)[1].strip())
                else:
                    correct_answers.append("")
            else:  # Subjective
                options.append([])
                # Subjective questions have no single "correct" answer
                correct_answers.append("")

    # Ensure we have the exact number of questions requested
    if len(questions) < num_questions:
        print(f"Warning: Only generated {len(questions)} questions instead of {num_questions}")
    elif len(questions) > num_questions:
        questions = questions[:num_questions]
        options = options[:num_questions]
        correct_answers = correct_answers[:num_questions]

    print(f"Successfully generated {len(questions)} questions")
    return questions, options, correct_answers

# Removed ensure_table_exists as MongoDB handles collection creation implicitly


def review_feedback():
    # Fetch feedback data from MongoDB
    db = create_connection()
    if db is None:
        return []

    try:
        feedback_data = list(db.feedback.find({}))
        # Convert ObjectId to string for compatibility with DataFrame
        for item in feedback_data:
            if '_id' in item:
                item['id'] = str(item['_id'])
                del item['_id']  # Remove ObjectId if not needed for display
            if 'question_bank_id' in item and isinstance(item['question_bank_id'], ObjectId):
                item['question_bank_id'] = str(item['question_bank_id'])
        return feedback_data
    except OperationFailure as e:
        st.error(f"Error retrieving feedback: {e}")
        return []


def analyze_sentiment(text):
    sia = SentimentIntensityAnalyzer()
    sentiment_score = sia.polarity_scores(text)['compound']
    if sentiment_score > 0.05:
        return 'Positive'
    elif sentiment_score < -0.05:
        return 'Negative'
    else:
        return 'Neutral'


def display_questions(questions, options, correct_answers):
    for i, question in enumerate(questions):
        st.write(question)
        if options[i]:
            st.write("Options:")
            for j, option in enumerate(options[i]):
                st.write(f"{chr(65+j)}) {option}")
        st.write(f"Correct Answer: {correct_answers[i]}")
        st.write("")
# User Authentication Functions


def login_user(username, password):
    db = create_connection()
    if db is None:
        return None

    user = db.users.find_one({"username": username})

    if user and check_password_hash(user['password'], password):
        # Convert ObjectId to string for session state
        user['_id'] = str(user['_id'])
        return user
    return None


def register_user(email, username, password, role):
    if db is None:
        return False

    # Basic validation
    if not re.match(r'^[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+$', email):
        st.error("Invalid email format.")
        return False
    if db.users.find_one({"username": username}):
        st.error("Username already exists.")
        return False

    hashed_password = generate_password_hash(password)
    user_data = {
        "email": email,
        "username": username,
        "password": hashed_password,
        "role": role,
    }
    # NEW: Set assigned_trainer to None for new employees
    if role == "Employee":
        user_data["assigned_trainer"] = None

    try:
        db.users.insert_one(user_data)
        # NEW: Notify trainers of new employee registration
        if role == "Employee":
            send_notification(
                recipient_role="trainer",
                message=f"New employee '{username}' has registered and needs assignment.",
                username=None  # General notification for all trainers
            )
        return True
    except OperationFailure as e:
        st.error(f"Database error during registration: {e}")
        return False
# Administrator Functions


def get_system_stats():
    db = create_connection()
    if db is None:
        return None

    stats = {}
    try:
        # Get count of documents in each collection
        stats['users'] = db.users.count_documents({})
        stats['question_banks'] = db.question_banks.count_documents({})
        stats['learning_plans'] = db.learning_plans.count_documents({})
        stats['feedback'] = db.feedback.count_documents({})
        return stats
    except OperationFailure as e:
        st.error(f"Error retrieving system stats: {e}")
        return None

# NEW: Fetches learning plans specifically for the logged-in user.


def get_user_learning_plans(username):
    """
    Retrieves all learning plans assigned to a specific user by joining
    learning_plans and question_banks collections.
    """
    db = create_connection()
    if db is None:
        return []
    try:
        pipeline = [
            {"$match": {"username": username}},
            {"$lookup": {
                "from": "question_banks",
                "localField": "question_bank_id",
                "foreignField": "_id",
                "as": "qb_details"
            }},
            {"$unwind": "$qb_details"},
            {"$project": {
                "_id": 1,  # The learning plan's unique ID
                "status": 1,
                "start_date": 1,
                "end_date": 1,
                "estimated_time": 1,
                "question_bank_id": "$qb_details._id",
                "technology": "$qb_details.technology",
                "difficulty": "$qb_details.difficulty"
            }}
        ]
        learning_plans = list(db.learning_plans.aggregate(pipeline))
        for lp in learning_plans:
            lp['_id'] = str(lp['_id'])
            lp['question_bank_id'] = str(lp['question_bank_id'])
        return learning_plans
    except OperationFailure as e:
        st.error(f"Database error retrieving learning plans: {e}")
        return []


def get_all_users():
    db = create_connection()
    if db is None:
        return []

    try:
        users_cursor = db.users.find(
            {}, {"username": 1, "email": 1, "role": 1, "_id": 0})
        users = list(users_cursor)
        return users
    except OperationFailure as e:
        st.error(f"Error retrieving all users: {e}")
        return []


def update_user_role(username, new_role):
    db = create_connection()
    if db is None:
        return False

    try:
        user_exists = db.users.find_one({"username": username})
        if not user_exists:
            st.error("User does not exist.")
            return False

        result = db.users.update_one({"username": username}, {
                                     "$set": {"role": new_role}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# Trainer Functions


def upload_curriculum(technology, topics, content, trainer_username):
    """Saves curriculum scoped to a specific trainer."""
    if db is None:
        return False
    try:
        curriculum_doc = {
            "trainer_username": trainer_username,
            "technology": technology,
            "topics": topics,
            "content": content,
            "created_at": datetime.now()
        }
        db.curriculum.update_one(
            {"trainer_username": trainer_username, "technology": technology},
            {"$set": curriculum_doc},
            upsert=True
        )
        return True
    except OperationFailure as e:
        st.error(f"Database error: {e}")
        return False


def get_curriculum_text(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one(
            {"technology": technology}, {"topics": 1, "_id": 0})
        if result:
            return result.get('topics')
        else:
            st.error(
                f"No curriculum content found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def save_question_bank(technology, topics, questions, difficulty, correct_answers, question_type, options, trainer_username):
    """Saves a question bank scoped to a specific trainer."""
    if db is None:
        return None
    try:
        qb_doc = {
            "trainer_username": trainer_username,
            "technology": technology,
            "topics": topics,
            "questions": questions,
            "difficulty": difficulty,
            "question_type": question_type,
            "options": options,
            "created_at": datetime.now()
        }
        result_qb = db.question_banks.insert_one(qb_doc)
        question_bank_id = result_qb.inserted_id

        answer_doc = {
            "question_bank_id": question_bank_id,
            "answer_data": correct_answers
        }
        db.question_answers.insert_one(answer_doc)
        return str(question_bank_id)
    except OperationFailure as e:
        st.error(f"MongoDB operation error: {e}")
        return None


def get_trainer_question_banks(trainer_username):
    """Fetches all question banks for a specific trainer."""
    if db is None:
        return []
    try:
        return list(db.question_banks.find({"trainer_username": trainer_username}))
    except OperationFailure as e:
        st.error(f"Database error: {e}")
        return []


def get_trainer_curricula(trainer_username):
    """Fetches all curricula for a specific trainer."""
    if db is None:
        return []
    try:
        return list(db.curriculum.find({"trainer_username": trainer_username}))
    except OperationFailure as e:
        st.error(f"Database error: {e}")
        return []


def get_topics_for_technology(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one(
            {"technology": technology}, {"topics": 1, "_id": 0})
        if result and 'topics' in result:
            return result['topics'].split(',')
        else:
            st.error(f"No topics found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def get_all_curricula():
    db = create_connection()
    if db is None:
        return None

    try:
        # Include _id for potential future use
        curricula_cursor = db.curriculum.find(
            {}, {"technology": 1, "topics": 1})
        curricula = []
        for doc in curricula_cursor:
            doc['id'] = str(doc['_id'])  # Add string version of _id
            curricula.append(doc)
        return curricula
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def get_all_question_banks():
    db = create_connection()
    if db is None:
        return None

    try:
        # Fetch all documents from the question_banks collection
        question_banks_cursor = db.question_banks.find({})
        question_banks = list(question_banks_cursor)
        return question_banks
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

# Employee Functions


# def get_learning_plan(username):
#     db = create_connection()
#     if db is None:
#         return None

#     learning_plan = db.learning_plans.find_one({"username": username})
#     return learning_plan


def submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
    sentiment = analyze_sentiment(feedback_text)  # Optional: Analyze sentiment
    db = create_connection()
    if db is None:
        return False

    try:
        feedback_doc = {
            "username": username,
            # Store as ObjectId
            "question_bank_id": ObjectId(question_bank_id) if question_bank_id else None,
            "feedback_text": feedback_text,
            "rating": rating,
            "sentiment": sentiment,
            "feedback_type": feedback_type,
            "created_at": datetime.now()
        }
        db.feedback.insert_one(feedback_doc)

        # Prepare notification message
        feedback_summary = f"New feedback received from {username}. Type: {feedback_type}. Rating: {rating}. Feedback: {feedback_text}"

        # Send notifications based on feedback type
        if feedback_type in ["User  Experience"] and rating >= 3:
            send_notification("admin", feedback_summary, username)
        elif feedback_type in ["Question Bank Feedback", "Assessment Feedback"]:
            send_notification("admin", feedback_summary, username)
            send_notification("trainer", feedback_summary, username)

        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


# def take_assessment():
#     st.subheader("Take Assessment")
#     question_banks = get_all_question_banks()
#     if not question_banks:
#         st.info("No question banks available yet.")
#     else:
#         selected_qb = st.selectbox(
#             "Select Question Bank",
#             options=[
#                 (str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
#             format_func=lambda x: f"ID: {x[0]} - {x[1]}",
#             key="take_assessment_qb_select"
#         )

#         if selected_qb:
#             qb_id_str, _ = selected_qb
#             qb_id = ObjectId(qb_id_str)  # Convert to ObjectId
#             qb_details = next(
#                 (qb for qb in question_banks if qb['_id'] == qb_id), None)
#             if qb_details:
#                 questions = qb_details.get('questions', '').split('\n')
#                 options = qb_details.get('options', '').split(
#                     '###') if qb_details.get('options') else []
#                 correct_answers_str = get_correct_answers(
#                     qb_id)  # This returns a list of strings

#                 question_type = qb_details.get('question_type')

#                 score = 0
#                 user_answers = []
#                 for i, question in enumerate(questions):
#                     if not question.strip():  # Skip empty questions
#                         continue

#                     st.write(f"**Q{i+1}:** {question.strip()}")

#                     if question_type == "multiple-choice" and options:
#                         # Assuming options are stored as 'Opt1###Opt2###Opt3###Opt4' per question
#                         # Need to parse options for the current question correctly
#                         # For now, let's assume options are structured to match questions
#                         if i < len(options):  # Ensure we have options for this question
#                             current_options = options[i].split('###')
#                             if current_options:
#                                 answer = st.radio(
#                                     "Select an option", current_options, key=f"question_{i}")
#                                 user_answers.append(answer)
#                             else:
#                                 # No options, no answer
#                                 user_answers.append("")
#                         else:
#                             # No options for this question
#                             user_answers.append("")
#                     elif question_type == "fill-in-the-blank":
#                         answer = st.text_input(
#                             "Enter your answer", key=f"question_{i}")
#                         user_answers.append(answer)
#                     elif question_type == "subjective":
#                         answer = st.text_area(
#                             "Enter your answer", key=f"question_{i}")
#                         user_answers.append(answer)
#                     else:
#                         # For cases where question_type is not set or options are missing
#                         user_answers.append("")

#                 if st.button("Submit"):
#                     total_questions_answered = 0
#                     correct_answers_count = 0
#                     for i, user_answer in enumerate(user_answers):
#                         if i < len(correct_answers_str) and user_answer.strip().lower() == correct_answers_str[i].strip().lower():
#                             st.success(f"Q{i+1}: Correct!")
#                             correct_answers_count += 1
#                         elif i < len(correct_answers_str):
#                             st.error(
#                                 f"Q{i+1}: Incorrect. Correct answer: {correct_answers_str[i].strip()}")
#                         total_questions_answered += 1

#                     st.write(
#                         f"Your score is {correct_answers_count} out of {total_questions_answered}")

#                     # Save the assessment result
#                     save_assessment_result(
#                         st.session_state.user['username'], qb_id, correct_answers_count)


# def get_available_question_banks(username):
#     db = create_connection()
#     if db is None:
#         return []

#     try:
#         # Get technologies from learning plans for the user
#         learning_plans_cursor = db.learning_plans.find(
#             {"username": username}, {"technology": 1, "_id": 0})
#         technologies = [lp['technology']
#                         for lp in learning_plans_cursor if 'technology' in lp]

#         # Get IDs of question banks already completed by the user
#         completed_assessments_cursor = db.assessments.find(
#             {"username": username}, {"question_bank_id": 1, "_id": 0})
#         completed_qb_ids = [a['question_bank_id']
#                             for a in completed_assessments_cursor if 'question_bank_id' in a]

#         # Find question banks matching learning plan technologies and not yet completed
#         query = {
#             "technology": {"$in": technologies},
#             "_id": {"$nin": completed_qb_ids}
#         }
#         question_banks_cursor = db.question_banks.find(
#             query, {"technology": 1, "topics": 1})

#         # Convert ObjectId to string for 'id' field in the returned dictionary
#         question_banks = []
#         for qb in question_banks_cursor:
#             qb['id'] = str(qb['_id'])
#             question_banks.append(qb)

#         return question_banks
#     except OperationFailure as e:
#         st.error(f"Error retrieving available question banks: {e}")
#         return []


def get_completed_assessments(username):
    db = create_connection()
    if db is None:
        return []

    try:
        # Fetch assessments for the user
        assessments_cursor = db.assessments.find(
            {"username": username}).sort("completed_at", -1)
        completed_assessments = []

        for assessment in assessments_cursor:
            # Fetch corresponding question bank details
            qb_details = db.question_banks.find_one(
                {"_id": assessment['question_bank_id']})

            if qb_details:
                total_questions = len(qb_details.get('questions', '').split(
                    '\n')) if qb_details.get('questions') else 0

                percentage = round(
                    (assessment['score'] / total_questions) * 100, 2) if total_questions > 0 else 0.0

                completed_assessments.append({
                    'id': str(assessment['_id']),  # Convert ObjectId to string
                    # Convert ObjectId to string
                    'question_bank_id': str(assessment['question_bank_id']),
                    'technology': qb_details.get('technology'),
                    'difficulty': qb_details.get('difficulty'),
                    'score': assessment.get('score'),
                    'question_type': qb_details.get('question_type'),
                    'completed_at': assessment.get('completed_at').strftime('%Y-%m-%d %H:%M:%S') if assessment.get('completed_at') else 'N/A',
                    'total_questions': total_questions,
                    'percentage': percentage
                })
        return completed_assessments
    except OperationFailure as e:
        st.error(f"Error retrieving completed assessments: {e}")
        return []


def admin_dashboard():

    with st.sidebar:
        # Create a sidebar for navigation using option_menu
        selected_tab = option_menu(
            menu_title="Admin Dashboard",  # required
            options=["System Stats",
                     "User  Management",
                     "Reports",
                     "Employee Performance"],  # required
            icons=["bar-chart", "people", "file-earmark-text",
                   "person-check"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="vertical",
        )

    # Display the selected tab content
    if selected_tab == "System Stats":
        st.subheader("System Statistics 📊")
        stats = get_system_stats()
        if stats:
            st.write(f"Total Users: {stats.get('users', 0)}")
            st.write(f"Total Question Banks: {stats.get('question_banks', 0)}")
            st.write(f"Total Learning Plans: {stats.get('learning_plans', 0)}")
            st.write(f"Total Feedback Entries: {stats.get('feedback', 0)}")

            # Add system details
            st.subheader("System Details 🖥️")
            st.write(f"Operating System: {platform.system()}")
            st.write(f"Platform: {platform.platform()}")
            st.write(f"Processor: {platform.processor()}")
            st.write(f"Python Version: {platform.python_version()}")
        else:
            st.error("Failed to retrieve system statistics")

    elif selected_tab == "User  Management":
        st.subheader("User  Management 👤")

        # Check if users are already loaded in session state
        if 'users' not in st.session_state:
            st.session_state.users = get_all_users()  # Load users for the first time

        # Create a layout for the refresh button
        col1, col2 = st.columns([4, 1])  # Adjust column widths

        with col1:
            st.write("")  # Empty space for alignment

        with col2:
            # Refresh button with an icon
            if st.button("🔄", key="refresh_users", help="Refresh User List"):
                st.session_state.users = get_all_users()  # Refresh the user list

        users = st.session_state.users  # Use the loaded users from session state

        if users:
            # Display the user table
            user_table = []
            for user in users:
                user_table.append({
                    'Username': user.get('username'),
                    'Email': user.get('email'),
                    'Role': user.get('role')
                })

            user_df = pd.DataFrame(user_table)
            st.table(user_df)

            # CSS to inject for compact layout
            st.markdown("""
                <style>
                .stSelectbox {
                    margin-bottom: 0px;
                }
                .stButton {
                    display: inline-block;
                    margin-right: 10px;
                }
                .user-row {
                    margin-bottom: 10px;
                }
                </style>
            """, unsafe_allow_html=True)

            # Search feature
            st.subheader("Search User")
            search_username = st.text_input("Enter username to search:", "")

            if search_username:
                filtered_users = [user for user in users if search_username.lower(
                ) in user.get('username', '').lower()]
            else:
                filtered_users = users

            # Display users (filtered or all)
            if filtered_users:
                for user in filtered_users:
                    with st.container():
                        cols = st.columns([2, 3, 3])  # Adjusted column widths

                        with cols[0]:  # Username column
                            st.write(user.get('username'))

                        with cols[1]:  # Role selection column
                            new_role = st.selectbox(
                                f"New Role for {user.get('username')}",
                                ["None", "Administrator", "Trainer", "Employee"],
                                key=f"new_role_{user.get('username')}",
                                label_visibility="collapsed"  # Hides the label
                            )

                        with cols[2]:  # Buttons column
                            # Split the column for buttons
                            c1, c2 = st.columns([1, 1])
                            with c1:
                                if st.button("Update Role", key=f"update_role_button_{user.get('username')}", use_container_width=True):
                                    if update_user_role(user.get('username'), new_role):
                                        st.success(
                                            f"Role updated for {user.get('username')}")
                                        # Refresh the user list after update
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to update role")

                            with c2:
                                if st.button("Remove User", key=f"remove_user_button_{user.get('username')}", use_container_width=True):
                                    if remove_user(user.get('username')):
                                        st.success(
                                            f"User {user.get('username')} removed successfully")
                                        # Refresh the user list after removal
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to remove user")
            else:
                st.info("No users found with that username.")
        else:
            st.info("No users available.")

    elif selected_tab == "Reports":
        st.subheader("Generate Reports 🔍")
        report_type = st.selectbox("Select Report Type",
                                   ["User Activity", "Question Bank Usage",
                                    "Feedback Summary", "Sentiment Analysis",
                                    "Employee Performance"])

        if st.button("Generate Report"):
            if report_type == "User Activity":
                user_activity_report()
            elif report_type == "Question Bank Usage":
                question_bank_usage_report()
            elif report_type == "Feedback Summary":
                feedback_summary_report()
            elif report_type == "Sentiment Analysis":
                sentiment_analysis_report()
            elif report_type == "Employee Performance":
                employee_performance_report()

    elif selected_tab == "Employee Performance":
        st.subheader("Employee Performance 🎯")
        employees = get_all_users()

        if employees:
            selected_employee = st.selectbox(
                "Select Employee",
                options=[employee['username'] for employee in employees],
                key="employee_performance_select"
            )

            if selected_employee:
                # Fetch assessment results for the selected employee
                assessment_results = get_assessment_results(selected_employee)
                if assessment_results:
                    # Prepare data for the table
                    performance_data = []
                    for result in assessment_results:
                        performance_data.append({
                            # Convert ObjectId to string
                            'Question Bank ID': str(result['question_bank_id']),
                            'Score': result['score'],
                            'Completed At': result['completed_at']
                        })

                    # Convert to DataFrame for better visualization
                    performance_df = pd.DataFrame(performance_data)

                    # Display summary metrics
                    st.subheader(f"Summary Statistics for {selected_employee}")
                    total_assessments = len(performance_df)
                    avg_score = performance_df['Score'].mean(
                    ) if total_assessments > 0 else 0
                    best_score = performance_df['Score'].max(
                    ) if total_assessments > 0 else 0

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Assessments", total_assessments)
                    with col2:
                        st.metric("Average Score", f"{avg_score:.1f}")
                    with col3:
                        st.metric("Best Score", best_score)

                    # Display the performance data in a styled table
                    st.write(f"Performance Data for {selected_employee}:")
                    st.dataframe(performance_df.style.highlight_max(
                        axis=0))  # Highlight max scores

                    # Convert 'Completed At' column to datetime for sorting
                    performance_df['Completed At'] = pd.to_datetime(
                        performance_df['Completed At'])
                    performance_df.sort_values('Completed At', inplace=True)

                    # Create visualizations
                    st.subheader("Performance Over Time")

                    # Line chart for scores over time
                    fig_line = px.line(performance_df, x='Completed At', y='Score',
                                       title='Score Over Time', markers=True)
                    st.plotly_chart(fig_line)

                    # Bar chart for scores by question bank
                    fig_bar = px.bar(performance_df, x='Question Bank ID', y='Score',
                                     title='Scores by Question Bank', text='Score')
                    st.plotly_chart(fig_bar)

                    # Convert figures to HTML format for download
                    fig_line_html = fig_line.to_html(full_html=False)
                    fig_bar_html = fig_bar.to_html(full_html=False)

                    # Provide download buttons
                    st.download_button(label="Download Line Chart as HTML", data=fig_line_html,
                                       file_name=f"{selected_employee}_performance_over_time.html", mime="text/html")

                    st.download_button(label="Download Bar Chart as HTML", data=fig_bar_html,
                                       file_name=f"{selected_employee}_score_by_question_bank.html", mime="text/html")

                    st.download_button(label="Download Performance Data as CSV",
                                       data=performance_df.to_csv(index=False),
                                       file_name=f"{selected_employee}_performance.csv", mime="text/csv")

                else:
                    st.info("No assessment results available for this employee.")
        else:
            st.info("No employees available.")

    # Display content based on the selected option
        if selected_tab == "Generate Questions":
            # Horizontal menu for question generation methods
            question_generation_method = option_menu(
                menu_title=None,  # required
                options=["Generate Questions by Topic",
                         "Generate Questions from Prompt"],  # required
                icons=["book", "pencil"],  # optional
                menu_icon="cast",  # optional
                default_index=0,  # optional
                orientation="horizontal",
            )

        if question_generation_method == "Generate Questions by Topic":
            st.subheader("Generate Questions by Topic 🚀")
            topic_name = st.text_input(
                "Enter Topic Name", key="topic_input")  # Input for topic name
            num_questions = st.number_input("Number of Questions to Generate", min_value=1,
                                            value=5, key="num_questions_input")  # Input for number of questions

            # Dropdown for selecting question type
            question_type = st.selectbox("Select Question Type", [
                                         "Multiple Choice", "Subjective", "Fill in the Blanks"], key="question_type_select")

            if st.button("Generate Questions", key="generate_topic_questions_button"):
                if topic_name:
                    try:
                        prompt = f"Generate {num_questions} {question_type.lower()} questions based on the topic: {topic_name}."
                        generated_questions = model.generate_content(
                            prompt)  # Using the model to generate questions
                        questions_text = generated_questions.text.strip()  # Extracting the generated text

                        # Displaying the generated questions
                        st.write("Generated Questions:")
                        questions_list = questions_text.split('\n')
                        selected_questions = []

                        for i, question in enumerate(questions_list, 1):
                            if question.strip():  # Only show non-empty questions
                                # Display each question with a checkbox
                                if st.checkbox(f"Question {i}: {question.strip()}", value=True):
                                    selected_questions.append(question.strip())

                        # Store generated questions in session state
                        if 'history' not in st.session_state:
                            st.session_state.history = []
                        st.session_state.history.append(
                            {"topic": topic_name, "questions": selected_questions})

                        st.session_state.generated_questions = selected_questions
                        st.session_state.topic_name = topic_name
                        st.success(
                            f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                    except Exception as e:
                        st.error(f"Error generating questions: {e}")
                else:
                    st.error("Please enter a topic name.")

            # New feature: Add Questions to Question Bank
            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")

                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [
                    f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox(
                    "Select Question Bank", options=qb_options)

                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input(
                            "Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", [
                                                         "Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(
                                new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(
                                    f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split(
                            # Convert to ObjectId
                            '-')[0].split(':')[1].strip())
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(
                                f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error(
                                "Failed to add questions to the selected question bank.")

                    # Clear the generated questions from session state
                    del st.session_state.generated_questions
                    del st.session_state.topic_name

        elif question_generation_method == "Generate Questions from Prompt":
            st.subheader("Generate Questions from Prompt ✍️")
            topic_name = st.text_input("Enter Topic Name")
            prompt = st.text_area("Enter a paragraph to generate questions")
            question_type = st.selectbox("Select Question Type", [
                                         "Multiple Choice", "Subjective", "Fill in the Blanks"])
            difficulty_level = st.selectbox("Select Difficulty Level", [
                                            "Easy", "Medium", "Hard"])
            num_questions = st.number_input(
                "Number of Questions to Generate", min_value=1, value=10)

            if st.button("Generate Questions"):
                generated_questions = generate_questions_from_prompt(
                    prompt, question_type, difficulty_level, num_questions, topic_name)
                if generated_questions:
                    st.write("Generated Questions:")
                    selected_questions = []
                    for i, question in enumerate(generated_questions):
                        if st.checkbox(f"Question {i+1}", value=True):
                            selected_questions.append(question)
                        st.write(question)

                    # Store generated questions in session state
                    if 'history' not in st.session_state:
                        st.session_state.history = []
                    st.session_state.history.append(
                        {"topic": topic_name, "questions": selected_questions})

                    st.session_state.generated_questions = selected_questions
                    st.session_state.topic_name = topic_name
                    st.success(
                        f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                else:
                    st.error("Failed to generate questions")

            # New feature: View Generated Questions History
            st.subheader("View Generated Questions History")
            topic_questions = get_generated_questions_history()  # Retrieve the history

            if topic_questions:
                selected_topic = st.selectbox(
                    "Select Topic", options=list(topic_questions.keys()))
                if selected_topic:
                    questions = topic_questions[selected_topic]
                    st.write(
                        f"Generated Questions for Topic: {selected_topic}")
                    for i, question in enumerate(questions, 1):
                        st.write(f"{i}. {question}")
            else:
                st.info("No generated questions history available.")

            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")

                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [
                    f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox(
                    "Select Question Bank", options=qb_options)

                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input(
                            "Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", [
                                                         "Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(
                                new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(
                                    f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split(
                            # Convert to ObjectId
                            '-')[0].split(':')[1].strip())
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(
                                f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error(
                                "Failed to add questions to the selected question bank.")

        # Clear the generated questions from session state
                        del st.session_state.generated_questions
                        del st.session_state.topic_name
        # Display content based on the selected option

    elif selected_tab == "Chatbot":
        # Display chatbot interface at the top
        st.subheader("Chat with the AI Trainer 🤖")

        # Initialize the message state if not exists
        if "msg" not in st.session_state:
            st.session_state.msg = ""

        # Create a container for the chat interface
        chat_container = st.container()

        # Define avatars
        user_avatar = "https://static.vecteezy.com/system/resources/previews/009/664/418/non_2x/people-user-team-transparent-free-png.png"
        ai_avatar = "https://thumbs.dreamstime.com/b/chatbot-logo-messenger-ai-robot-icon-vector-illustration-277900892.jpg"

        def clear_text():
            st.session_state.msg = st.session_state.user_input
            st.session_state.user_input = ""

        with chat_container:
            # Display chat messages
            for chat in st.session_state.chat_history:
                if chat['role'] == 'assistant':
                    # Chatbot message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;color:black'>"
                        f"<img src='{ai_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-right: 10px;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>AI:</strong> {chat['content']}</div></div>",
                        unsafe_allow_html=True
                    )
                else:
                    # User message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0; justify-content: flex-end;color:black'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {chat['content']}</div>"
                        f"<img src='{user_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-left: 10px;'>"
                        f"</div>",
                        unsafe_allow_html=True
                    )

            # Input field for user to enter a prompt
            st.text_input("Type your message here...", key="user_input",
                          placeholder="Type a message...", on_change=clear_text)

            if st.session_state.msg:  # Only process if there's a message
                # Append user input to chat history
                st.session_state.chat_history.append(
                    {"role": "user", "content": st.session_state.msg})

                try:
                    # Generate AI response
                    prompt = f"You are an AI assistant for trainers. Respond to the following message: {st.session_state.msg}"
                    response = model.generate_content(prompt)

                    # Handle the response properly for Gemini model
                    if hasattr(response, 'parts'):
                        ai_response = ''.join(
                            part.text for part in response.parts)
                    else:
                        ai_response = response.candidates[0].content.parts[0].text

                    # Append AI response to chat history
                    st.session_state.chat_history.append(
                        {"role": "assistant", "content": ai_response})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    ai_response = "I apologize, but I encountered an error. Please try again."
                    st.session_state.chat_history.append(
                        {"role": "assistant", "content": ai_response})

                # Clear the message state
                st.session_state.msg = ""

                # Rerun the app to display the new messages
                st.rerun()

    # Get notifications for trainer
    notifications = get_notifications("trainer", None)
    # Display notifications in the sidebar
    display_notifications(notifications, "trainer")

    if notifications:
        st.sidebar.write("Notifications:")
        for notification in notifications:
            st.sidebar.write(notification['message'])
    else:
        st.sidebar.write("No notifications available.")


# Download NLTK data
nltk.download('vader_lexicon')


def ascii_to_string(ascii_list):
    """Convert a list of ASCII values to a string."""
    return ''.join(chr(num) for num in ascii_list)

    # Example usage when retrieving topics
    curricula = get_all_curricula()
    for curriculum in curricula:
        # Assuming 'topics' is a list of ASCII values in the curricula
        # Check if topics is a list of ASCII values
        if isinstance(curriculum['topics'], list):
            curriculum['topics'] = ascii_to_string(
                curriculum['topics'])  # Convert to string


def format_timestamp(timestamp_str):
    """
    Format timestamp for display
    """
    if timestamp_str == 'N/A':
        return 'N/A'
    try:
        # Parse the timestamp string to datetime
        timestamp = datetime.datetime.strptime(
            timestamp_str, "%Y-%m-%d %H:%M:%S")
        # Format it to a more readable format
        return timestamp.strftime("%b %d, %Y %I:%M %p")
    except (ValueError, TypeError):
        return 'N/A'


# Utility functions
def extract_text_from_file(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    text = ""

    try:
        if file_extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_extension == '.docx':
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == '.txt':
            text = file.getvalue().decode('utf-8')
        elif file_extension in ['.ppt', '.pptx']:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        elif file_extension == '.csv':
            csv_data = pd.read_csv(file)
            text = csv_data.to_string(index=False)
        else:
            # For other file types, attempt to read as text
            try:
                text = file.getvalue().decode('utf-8')
            except UnicodeDecodeError:
                raise ValueError(
                    f"Unable to extract text from {file_extension} file.")
    except Exception as e:
        raise ValueError(f"Error processing {file_extension} file: {str(e)}")

    # Clean the extracted text
    cleaned_text = clean_text(text)
    return cleaned_text


def clean_text(text):
    # Remove non-printable characters and control characters
    # Keep only printable ASCII characters
    text = re.sub(r'[^\x20-\x7E]+', ' ', text)

    # Normalize whitespace
    # Replace multiple spaces with a single space
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()  # Trim leading and trailing whitespace

    return text


def create_new_question_bank(technology, difficulty, questions):
    db = create_connection()
    if db is None:
        return None

    try:
        questions_text = '\n'.join(questions)
        new_qb_doc = {
            "technology": technology,
            "difficulty": difficulty,
            "questions": questions_text,
            "created_at": datetime.now()
        }
        result = db.question_banks.insert_one(new_qb_doc)
        return str(result.inserted_id)  # Return the ObjectId as a string
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_questions(text, num_questions=5, question_type="multiple-choice"):
    if question_type == "multiple-choice":
        prompt = f"Generate {num_questions} multiple-choice questions based on the following text:\n\n{text}\n\nProvide the questions and options in the following format:\n\nQ1: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\nQ2: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\n..."
    elif question_type == "subjective":
        prompt = f"Generate {num_questions} subjective questions based on the following text:\n\n{text}\n\nProvide the questions in the following format:\n\nQ1: [Question]\n\nQ2: [Question]\n\n..."
    elif question_type == "fill-in-the-blank":
        prompt = f"Generate {num_questions} fill-in-the-blank questions based on the following text:\n\n{text}\n\nProvide the questions and correct answers in the following format:\n\nQ1: [Question]\nA: [Correct Answer]\n\nQ2: [Question]\nA: [Correct Answer]\n\n..."
    else:
        raise ValueError("Invalid question type")

    response = model.generate_content(prompt)
    generated_text = response.text

    questions = []
    options = []
    correct_answers = []

    lines = [line.strip()
             for line in generated_text.split('\n') if line.strip()]

    i = 0
    while i < len(lines):
        if lines[i].startswith('Q'):
            question = lines[i].split(': ', 1)[1]
            questions.append(question)
            if question_type == "multiple-choice":
                options_list = []
                correct_answer = None
                # Look for options immediately following the question
                for j in range(i + 1, len(lines)):
                    if lines[j].startswith(('A)', 'B)', 'C)', 'D)')):
                        option = lines[j].split(') ', 1)[1]
                        options_list.append(option)
                        # Assuming A is always the correct answer for simplicity in parsing
                        if lines[j].startswith('A)'):
                            correct_answer = option
                    else:
                        break  # Stop if a line doesn't start with an option letter
                options.append(options_list)
                correct_answers.append(correct_answer)
                i = j  # Move index to the line after the last option processed
            elif question_type == "fill-in-the-blank":
                if i+1 < len(lines) and lines[i+1].startswith('A:'):
                    options.append([lines[i+1].split(': ', 1)[1]])
                    correct_answers.append(lines[i+1].split(': ', 1)[1])
                    i += 2
                else:
                    options.append([""])
                    correct_answers.append("")
                    i += 1
            else:  # subjective
                options.append([])
                correct_answers.append("")
                i += 1
        else:
            i += 1

    return questions[:num_questions], options[:num_questions], correct_answers[:num_questions]

# Removed ensure_table_exists as MongoDB handles collection creation implicitly


def review_feedback():
    # Fetch feedback data from MongoDB
    db = create_connection()
    if db is None:
        return []

    try:
        feedback_data = list(db.feedback.find({}))
        # Convert ObjectId to string for compatibility with DataFrame
        for item in feedback_data:
            if '_id' in item:
                item['id'] = str(item['_id'])
                del item['_id']  # Remove ObjectId if not needed for display
            if 'question_bank_id' in item and isinstance(item['question_bank_id'], ObjectId):
                item['question_bank_id'] = str(item['question_bank_id'])
        return feedback_data
    except OperationFailure as e:
        st.error(f"Error retrieving feedback: {e}")
        return []


def analyze_sentiment(text):
    sia = SentimentIntensityAnalyzer()
    sentiment_score = sia.polarity_scores(text)['compound']
    if sentiment_score > 0.05:
        return 'Positive'
    elif sentiment_score < -0.05:
        return 'Negative'
    else:
        return 'Neutral'


def display_questions(questions, options, correct_answers):
    for i, question in enumerate(questions):
        st.write(question)
        if options[i]:
            st.write("Options:")
            for j, option in enumerate(options[i]):
                st.write(f"{chr(65+j)}) {option}")
        st.write(f"Correct Answer: {correct_answers[i]}")
        st.write("")
# User Authentication Functions


def login_user(username, password):
    db = create_connection()
    if db is None:
        return None

    user = db.users.find_one({"username": username})

    if user and check_password_hash(user['password'], password):
        # Convert ObjectId to string for session state
        user['_id'] = str(user['_id'])
        return user
    return None


def register_user(email, username, password, role):
    # Validate email format
    email_regex = r'^[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+$'
    if not re.match(email_regex, email):
        st.error("Invalid email format. Please enter a valid email address.")
        return False

    db = create_connection()
    if db is None:
        return False

    existing_user = db.users.find_one({"username": username})

    if existing_user:
        return False

    hashed_password = generate_password_hash(password)
    user_data = {
        "email": email,
        "username": username,
        "password": hashed_password,
        "role": role
    }
    try:
        db.users.insert_one(user_data)
        return True
    except OperationFailure as e:
        st.error(f"Database error during registration: {e}")
        return False
# Administrator Functions


def get_system_stats():
    db = create_connection()
    if db is None:
        return None

    stats = {}
    try:
        # Get count of documents in each collection
        stats['users'] = db.users.count_documents({})
        stats['question_banks'] = db.question_banks.count_documents({})
        stats['learning_plans'] = db.learning_plans.count_documents({})
        stats['feedback'] = db.feedback.count_documents({})
        return stats
    except OperationFailure as e:
        st.error(f"Error retrieving system stats: {e}")
        return None


def get_all_users():
    db = create_connection()
    if db is None:
        return []

    try:
        users_cursor = db.users.find(
            {}, {"username": 1, "email": 1, "role": 1, "_id": 0})
        users = list(users_cursor)
        return users
    except OperationFailure as e:
        st.error(f"Error retrieving all users: {e}")
        return []


def update_user_role(username, new_role):
    db = create_connection()
    if db is None:
        return False

    try:
        user_exists = db.users.find_one({"username": username})
        if not user_exists:
            st.error("User does not exist.")
            return False

        result = db.users.update_one({"username": username}, {
                                     "$set": {"role": new_role}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# Trainer Functions


def upload_curriculum(technology, topics, content, trainer_username):
    """Saves curriculum scoped to a specific trainer."""
    db = create_connection()
    if db is None:
        return False
        
    try:
        curriculum_doc = {
            "trainer_username": trainer_username, # <-- Saves the trainer's name
            "technology": technology,
            "topics": topics,
            "content": content,
            "created_at": datetime.now()
        }
        # Use upsert to either create a new curriculum or update an existing one
        db.curriculum.update_one(
            {"trainer_username": trainer_username, "technology": technology},
            {"$set": curriculum_doc},
            upsert=True
        )
        return True
    except OperationFailure as e:
        st.error(f"Database error: {e}")
        return False


def get_curriculum_text(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one(
            {"technology": technology}, {"topics": 1, "_id": 0})
        if result:
            return result.get('topics')
        else:
            st.error(
                f"No curriculum content found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None



def get_topics_for_technology(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one(
            {"technology": technology}, {"topics": 1, "_id": 0})
        if result and 'topics' in result:
            return result['topics'].split(',')
        else:
            st.error(f"No topics found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def get_all_curricula():
    db = create_connection()
    if db is None:
        return None

    try:
        # Include _id for potential future use
        curricula_cursor = db.curriculum.find(
            {}, {"technology": 1, "topics": 1})
        curricula = []
        for doc in curricula_cursor:
            doc['id'] = str(doc['_id'])  # Add string version of _id
            curricula.append(doc)
        return curricula
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def get_all_question_banks():
    db = create_connection()
    if db is None:
        return None

    try:
        # Fetch all documents from the question_banks collection
        question_banks_cursor = db.question_banks.find({})
        question_banks = list(question_banks_cursor)
        return question_banks
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

# Employee Functions


def get_learning_plan(username):
    db = create_connection()
    if db is None:
        return None

    learning_plan = db.learning_plans.find_one({"username": username})
    return learning_plan


def submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
    sentiment = analyze_sentiment(feedback_text)  # Optional: Analyze sentiment
    db = create_connection()
    if db is None:
        return False

    try:
        feedback_doc = {
            "username": username,
            # Store as ObjectId
            "question_bank_id": ObjectId(question_bank_id) if question_bank_id else None,
            "feedback_text": feedback_text,
            "rating": rating,
            "sentiment": sentiment,
            "feedback_type": feedback_type,
            "created_at": datetime.now()
        }
        db.feedback.insert_one(feedback_doc)

        # Prepare notification message
        feedback_summary = f"New feedback received from {username}. Type: {feedback_type}. Rating: {rating}. Feedback: {feedback_text}"

        # Send notifications based on feedback type
        if feedback_type in ["User  Experience"] and rating >= 3:
            send_notification("admin", feedback_summary, username)
        elif feedback_type in ["Question Bank Feedback", "Assessment Feedback"]:
            send_notification("admin", feedback_summary, username)
            send_notification("trainer", feedback_summary, username)

        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


# def take_assessment():
#     st.subheader("Take Assessment")
#     question_banks = get_all_question_banks()
#     if not question_banks:
#         st.info("No question banks available yet.")
#     else:
#         selected_qb = st.selectbox(
#             "Select Question Bank",
#             options=[
#                 (str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
#             format_func=lambda x: f"ID: {x[0]} - {x[1]}",
#             key="take_assessment_qb_select"
#         )

#         if selected_qb:
#             qb_id_str, _ = selected_qb
#             qb_id = ObjectId(qb_id_str)  # Convert to ObjectId
#             qb_details = next(
#                 (qb for qb in question_banks if qb['_id'] == qb_id), None)
#             if qb_details:
#                 questions = qb_details.get('questions', '').split('\n')
#                 options = qb_details.get('options', '').split(
#                     '###') if qb_details.get('options') else []
#                 correct_answers_str = get_correct_answers(
#                     qb_id)  # This returns a list of strings

#                 question_type = qb_details.get('question_type')

#                 score = 0
#                 user_answers = []
#                 for i, question in enumerate(questions):
#                     if not question.strip():  # Skip empty questions
#                         continue

#                     st.write(f"**Q{i+1}:** {question.strip()}")

#                     if question_type == "multiple-choice" and options:
#                         # Assuming options are stored as 'Opt1###Opt2###Opt3###Opt4' per question
#                         # Need to parse options for the current question correctly
#                         # For now, let's assume options are structured to match questions
#                         if i < len(options):  # Ensure we have options for this question
#                             current_options = options[i].split('###')
#                             if current_options:
#                                 answer = st.radio(
#                                     "Select an option", current_options, key=f"question_{i}")
#                                 user_answers.append(answer)
#                             else:
#                                 # No options, no answer
#                                 user_answers.append("")
#                         else:
#                             # No options for this question
#                             user_answers.append("")
#                     elif question_type == "fill-in-the-blank":
#                         answer = st.text_input(
#                             "Enter your answer", key=f"question_{i}")
#                         user_answers.append(answer)
#                     elif question_type == "subjective":
#                         answer = st.text_area(
#                             "Enter your answer", key=f"question_{i}")
#                         user_answers.append(answer)
#                     else:
#                         # For cases where question_type is not set or options are missing
#                         user_answers.append("")

#                 if st.button("Submit"):
#                     total_questions_answered = 0
#                     correct_answers_count = 0
#                     for i, user_answer in enumerate(user_answers):
#                         if i < len(correct_answers_str) and user_answer.strip().lower() == correct_answers_str[i].strip().lower():
#                             st.success(f"Q{i+1}: Correct!")
#                             correct_answers_count += 1
#                         elif i < len(correct_answers_str):
#                             st.error(
#                                 f"Q{i+1}: Incorrect. Correct answer: {correct_answers_str[i].strip()}")
#                         total_questions_answered += 1

#                     st.write(
#                         f"Your score is {correct_answers_count} out of {total_questions_answered}")

#                     # Save the assessment result
#                     save_assessment_result(
#                         st.session_state.user['username'], qb_id, correct_answers_count)


# def get_available_question_banks(username):
#     db = create_connection()
#     if db is None:
#         return []

#     try:
#         # Get technologies from learning plans for the user
#         learning_plans_cursor = db.learning_plans.find(
#             {"username": username}, {"technology": 1, "_id": 0})
#         technologies = [lp['technology']
#                         for lp in learning_plans_cursor if 'technology' in lp]

#         # Get IDs of question banks already completed by the user
#         completed_assessments_cursor = db.assessments.find(
#             {"username": username}, {"question_bank_id": 1, "_id": 0})
#         completed_qb_ids = [a['question_bank_id']
#                             for a in completed_assessments_cursor if 'question_bank_id' in a]

#         # Find question banks matching learning plan technologies and not yet completed
#         query = {
#             "technology": {"$in": technologies},
#             "_id": {"$nin": completed_qb_ids}
#         }
#         question_banks_cursor = db.question_banks.find(
#             query, {"technology": 1, "topics": 1})

#         # Convert ObjectId to string for 'id' field in the returned dictionary
#         question_banks = []
#         for qb in question_banks_cursor:
#             qb['id'] = str(qb['_id'])
#             question_banks.append(qb)

#         return question_banks
#     except OperationFailure as e:
#         st.error(f"Error retrieving available question banks: {e}")
#         return []


def get_completed_assessments(username):
    db = create_connection()
    if db is None:
        return []

    try:
        # Fetch assessments for the user
        assessments_cursor = db.assessments.find(
            {"username": username}).sort("completed_at", -1)
        completed_assessments = []

        for assessment in assessments_cursor:
            # Fetch corresponding question bank details
            qb_details = db.question_banks.find_one(
                {"_id": assessment['question_bank_id']})

            if qb_details:
                total_questions = len(qb_details.get('questions', '').split(
                    '\n')) if qb_details.get('questions') else 0

                percentage = round(
                    (assessment['score'] / total_questions) * 100, 2) if total_questions > 0 else 0.0

                completed_assessments.append({
                    'id': str(assessment['_id']),  # Convert ObjectId to string
                    # Convert ObjectId to string
                    'question_bank_id': str(assessment['question_bank_id']),
                    'technology': qb_details.get('technology'),
                    'difficulty': qb_details.get('difficulty'),
                    'score': assessment.get('score'),
                    'question_type': qb_details.get('question_type'),
                    'completed_at': assessment.get('completed_at').strftime('%Y-%m-%d %H:%M:%S') if assessment.get('completed_at') else 'N/A',
                    'total_questions': total_questions,
                    'percentage': percentage
                })
        return completed_assessments
    except OperationFailure as e:
        st.error(f"Error retrieving completed assessments: {e}")
        return []


def admin_dashboard():

    with st.sidebar:
        # Create a sidebar for navigation using option_menu
        selected_tab = option_menu(
            menu_title="Admin Dashboard",  # required
            options=["System Stats",
                     "User  Management",
                     "Reports",
                     "Employee Performance"],  # required
            icons=["bar-chart", "people", "file-earmark-text",
                   "person-check"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="vertical",
        )

    # Display the selected tab content
    if selected_tab == "System Stats":
        st.subheader("System Statistics 📊")
        stats = get_system_stats()
        if stats:
            st.write(f"Total Users: {stats.get('users', 0)}")
            st.write(f"Total Question Banks: {stats.get('question_banks', 0)}")
            st.write(f"Total Learning Plans: {stats.get('learning_plans', 0)}")
            st.write(f"Total Feedback Entries: {stats.get('feedback', 0)}")

            # Add system details
            st.subheader("System Details 🖥️")
            st.write(f"Operating System: {platform.system()}")
            st.write(f"Platform: {platform.platform()}")
            st.write(f"Processor: {platform.processor()}")
            st.write(f"Python Version: {platform.python_version()}")
        else:
            st.error("Failed to retrieve system statistics")

    elif selected_tab == "User  Management":
        st.subheader("User  Management �")

        # Check if users are already loaded in session state
        if 'users' not in st.session_state:
            st.session_state.users = get_all_users()  # Load users for the first time

        # Create a layout for the refresh button
        col1, col2 = st.columns([4, 1])  # Adjust column widths

        with col1:
            st.write("")  # Empty space for alignment

        with col2:
            # Refresh button with an icon
            if st.button("🔄", key="refresh_users", help="Refresh User List"):
                st.session_state.users = get_all_users()  # Refresh the user list

        users = st.session_state.users  # Use the loaded users from session state

        if users:
            # Display the user table
            user_table = []
            for user in users:
                user_table.append({
                    'Username': user.get('username'),
                    'Email': user.get('email'),
                    'Role': user.get('role')
                })

            user_df = pd.DataFrame(user_table)
            st.table(user_df)

            # CSS to inject for compact layout
            st.markdown("""
                <style>
                .stSelectbox {
                    margin-bottom: 0px;
                }
                .stButton {
                    display: inline-block;
                    margin-right: 10px;
                }
                .user-row {
                    margin-bottom: 10px;
                }
                </style>
            """, unsafe_allow_html=True)

            # Search feature
            st.subheader("Search User")
            search_username = st.text_input("Enter username to search:", "")

            if search_username:
                filtered_users = [user for user in users if search_username.lower(
                ) in user.get('username', '').lower()]
            else:
                filtered_users = users

            # Display users (filtered or all)
            if filtered_users:
                for user in filtered_users:
                    with st.container():
                        cols = st.columns([2, 3, 3])  # Adjusted column widths

                        with cols[0]:  # Username column
                            st.write(user.get('username'))

                        with cols[1]:  # Role selection column
                            new_role = st.selectbox(
                                f"New Role for {user.get('username')}",
                                ["None", "Administrator", "Trainer", "Employee"],
                                key=f"new_role_{user.get('username')}",
                                label_visibility="collapsed"  # Hides the label
                            )

                        with cols[2]:  # Buttons column
                            # Split the column for buttons
                            c1, c2 = st.columns([1, 1])
                            with c1:
                                if st.button("Update Role", key=f"update_role_button_{user.get('username')}", use_container_width=True):
                                    if update_user_role(user.get('username'), new_role):
                                        st.success(
                                            f"Role updated for {user.get('username')}")
                                        # Refresh the user list after update
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to update role")

                            with c2:
                                if st.button("Remove User", key=f"remove_user_button_{user.get('username')}", use_container_width=True):
                                    if remove_user(user.get('username')):
                                        st.success(
                                            f"User {user.get('username')} removed successfully")
                                        # Refresh the user list after removal
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to remove user")
            else:
                st.info("No users found with that username.")
        else:
            st.info("No users available.")

    elif selected_tab == "Reports":
        st.subheader("Generate Reports 🔍")
        report_type = st.selectbox("Select Report Type",
                                   ["User Activity", "Question Bank Usage",
                                    "Feedback Summary", "Sentiment Analysis",
                                    "Employee Performance"])

        if st.button("Generate Report"):
            if report_type == "User Activity":
                user_activity_report()
            elif report_type == "Question Bank Usage":
                question_bank_usage_report()
            elif report_type == "Feedback Summary":
                feedback_summary_report()
            elif report_type == "Sentiment Analysis":
                sentiment_analysis_report()
            elif report_type == "Employee Performance":
                employee_performance_report()

    elif selected_tab == "Employee Performance":
        st.subheader("Employee Performance 🎯")
        employees = get_all_users()

        if employees:
            selected_employee = st.selectbox(
                "Select Employee",
                options=[employee['username'] for employee in employees],
                key="employee_performance_select"
            )

            if selected_employee:
                # Fetch assessment results for the selected employee
                assessment_results = get_assessment_results(selected_employee)
                if assessment_results:
                    # Prepare data for the table
                    performance_data = []
                    for result in assessment_results:
                        performance_data.append({
                            # Convert ObjectId to string
                            'Question Bank ID': str(result['question_bank_id']),
                            'Score': result['score'],
                            'Completed At': result['completed_at']
                        })

                    # Convert to DataFrame for better visualization
                    performance_df = pd.DataFrame(performance_data)

                    # Display summary metrics
                    st.subheader(f"Summary Statistics for {selected_employee}")
                    total_assessments = len(performance_df)
                    avg_score = performance_df['Score'].mean(
                    ) if total_assessments > 0 else 0
                    best_score = performance_df['Score'].max(
                    ) if total_assessments > 0 else 0

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Assessments", total_assessments)
                    with col2:
                        st.metric("Average Score", f"{avg_score:.1f}")
                    with col3:
                        st.metric("Best Score", best_score)

                    # Display the performance data in a styled table
                    st.write(f"Performance Data for {selected_employee}:")
                    st.dataframe(performance_df.style.highlight_max(
                        axis=0))  # Highlight max scores

                    # Convert 'Completed At' column to datetime for sorting
                    performance_df['Completed At'] = pd.to_datetime(
                        performance_df['Completed At'])
                    performance_df.sort_values('Completed At', inplace=True)

                    # Create visualizations
                    st.subheader("Performance Over Time")

                    # Line chart for scores over time
                    fig_line = px.line(performance_df, x='Completed At', y='Score',
                                       title='Score Over Time', markers=True)
                    st.plotly_chart(fig_line)

                    # Bar chart for scores by question bank
                    fig_bar = px.bar(performance_df, x='Question Bank ID', y='Score',
                                     title='Scores by Question Bank', text='Score')
                    st.plotly_chart(fig_bar)

                    # Convert figures to HTML format for download
                    fig_line_html = fig_line.to_html(full_html=False)
                    fig_bar_html = fig_bar.to_html(full_html=False)

                    # Provide download buttons
                    st.download_button(label="Download Line Chart as HTML", data=fig_line_html,
                                       file_name=f"{selected_employee}_performance_over_time.html", mime="text/html")

                    st.download_button(label="Download Bar Chart as HTML", data=fig_bar_html,
                                       file_name=f"{selected_employee}_score_by_question_bank.html", mime="text/html")

                    st.download_button(label="Download Performance Data as CSV",
                                       data=performance_df.to_csv(index=False),
                                       file_name=f"{selected_employee}_performance.csv", mime="text/csv")

                else:
                    st.info("No assessment results available for this employee.")
        else:
            st.info("No employees available.")

    # Display notifications in the sidebar

    notifications = get_notifications(
        "admin", None)  # Get notifications for admin
    # Display notifications in the sidebar
    display_notifications(notifications, "admin")
    if notifications:
        st.sidebar.write("Notifications:")
        for notification in notifications:
            st.sidebar.write(notification['message'])
    else:
        st.sidebar.write("No notifications available.")


def employee_performance_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get employee performance
        pipeline = [
            {"$group": {
                "_id": "$username",
                "num_assessments": {"$sum": 1},
                "avg_score": {"$avg": "$score"}
            }},
            {"$project": {
                "username": "$_id",
                "num_assessments": 1,
                "avg_score": 1,
                "_id": 0
            }}
        ]
        results = list(db.assessments.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Create an interactive bar chart using Plotly
            fig = px.bar(df, x='username', y='avg_score', title='Employee Performance',
                         labels={'username': 'Employee',
                                 'avg_score': 'Average Score'},
                         color='avg_score', color_continuous_scale=px.colors.sequential.Viridis)

            # Display the chart in Streamlit
            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Provide a download button for the interactive HTML chart
            st.download_button(
                label="Download Chart as HTML",
                data=fig_html,
                file_name="employee_performance_chart.html",
                mime="text/html"
            )

            # Download CSV Data
            csv = df.to_csv(index=False)
            st.download_button(label="Download as CSV", data=csv,
                               file_name="employee_performance.csv", mime="text/csv")

        else:
            st.error("No employee performance data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def user_activity_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get user activity for employees
        pipeline = [
            {"$lookup": {
                "from": "users",
                "localField": "username",
                "foreignField": "username",
                "as": "user_info"
            }},
            {"$unwind": "$user_info"},
            {"$match": {"user_info.role": "Employee"}},
            {"$group": {
                "_id": "$username",
                "num_assessments": {"$sum": 1},
                "avg_score": {"$avg": "$score"}
            }},
            {"$project": {
                "username": "$_id",
                "num_assessments": 1,
                "avg_score": 1,
                "_id": 0
            }}
        ]
        results = list(db.assessments.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Pie Chart
            fig = px.pie(df, values='num_assessments',
                         names='username', title='User Activity Distribution')
            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html,
                               file_name="user_activity_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(
                index=False), file_name="user_activity.csv", mime="text/csv")

        else:
            st.error("No user activity data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def remove_user(username):
    db = create_connection()
    if db is None:
        return False

    try:
        # First, delete related records in the assessments collection
        db.assessments.delete_many({"username": username})

        # Now, delete the user from the users collection
        result = db.users.delete_one({"username": username})

        return result.deleted_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def question_bank_usage_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get question bank usage
        pipeline = [
            {"$lookup": {
                "from": "question_banks",
                "localField": "question_bank_id",
                "foreignField": "_id",
                "as": "qb_info"
            }},
            {"$unwind": "$qb_info"},
            {"$group": {
                "_id": "$qb_info.technology",
                "num_assessments": {"$sum": 1}
            }},
            {"$project": {
                "technology": "$_id",
                "num_assessments": 1,
                "_id": 0
            }}
        ]
        results = list(db.assessments.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Bar Chart
            fig = px.bar(df, x='technology', y='num_assessments', title='Question Bank Usage',
                         labels={'technology': 'Technology',
                                 'num_assessments': 'Number of Assessments'},
                         color='num_assessments', color_continuous_scale=px.colors.sequential.Plasma)

            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html,
                               file_name="question_bank_usage_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(
                index=False), file_name="question_bank_usage.csv", mime="text/csv")

        else:
            st.error("No question bank usage data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def sentiment_analysis_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get sentiment counts
        pipeline = [
            {"$group": {
                "_id": "$sentiment",
                "num_feedback": {"$sum": 1}
            }},
            {"$project": {
                "sentiment": "$_id",
                "num_feedback": 1,
                "_id": 0
            }}
        ]
        results = list(db.feedback.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Pie Chart
            fig = px.pie(df, values='num_feedback', names='sentiment', title='Sentiment Analysis',
                         color='sentiment', color_discrete_sequence=px.colors.qualitative.Set1)

            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html,
                               file_name="sentiment_analysis_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(
                index=False), file_name="sentiment_analysis.csv", mime="text/csv")

        else:
            st.error("No sentiment analysis data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def notify_trainer_of_bad_feedback(feedback_summary):
    send_notification("trainer", feedback_summary, "Admin Notification")


def feedback_summary_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get feedback summary
        pipeline = [
            {"$group": {
                "_id": "$feedback_type",
                "avg_rating": {"$avg": "$rating"},
                "num_feedback": {"$sum": 1}
            }},
            {"$project": {
                "feedback_type": "$_id",
                "avg_rating": 1,
                "num_feedback": 1,
                "_id": 0
            }}
        ]
        results = list(db.feedback.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Bar Chart
            fig = px.bar(df, x='feedback_type', y='num_feedback', title='Feedback Summary',
                         labels={'feedback_type': 'Feedback Type',
                                 'num_feedback': 'Number of Feedback'},
                         color='num_feedback', color_continuous_scale=px.colors.sequential.Viridis)

            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html,
                               file_name="feedback_summary_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(
                index=False), file_name="feedback_summary.csv", mime="text/csv")

        else:
            st.error("No feedback data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def get_bad_feedback():
    db = create_connection()
    if db is None:
        return []

    try:
        feedback_records = list(db.feedback.find({"rating": {"$lt": 3}}, {
                                "username": 1, "feedback_text": 1, "rating": 1, "_id": 0}))
        return feedback_records
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def feedback_form(username):
    st.subheader("Submit Feedback")

    # Select the type of feedback
    feedback_type = st.selectbox("Select Feedback Type:",
                                 options=["User  Experience", "Question Bank Feedback", "Assessment Feedback"])

    # Select a question bank for feedback if applicable
    question_banks = get_all_question_banks()
    selected_qb = None
    if feedback_type in ["Question Bank Feedback", "Assessment Feedback"] and question_banks:
        selected_qb = st.selectbox(
            "Select Question Bank for Feedback",
            options=[
                (str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
            format_func=lambda x: f"ID: {x[0]} - {x[1]}",
            key="feedback_qb_select"
        )

    feedback_text = st.text_area("Submit your feedback:", height=100)
    rating = st.slider("Rate your experience (1-5):", 1, 5, 3)

    if st.button("Submit Feedback"):
        question_bank_id = selected_qb[0] if selected_qb else None
        if submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
            st.success("Feedback submitted successfully!")
        else:
            st.error("Failed to submit feedback.")


def feedback_received(feedback):
    db = create_connection()
    if db is None:
        return False

    try:
        feedback_doc = {"feedback": feedback, "created_at": datetime.now()}
        db.feedback.insert_one(feedback_doc)

        print("Feedback received!")

        # Send notification to admin
        recipient_role = "admin"
        message = f"New feedback received: {feedback}"
        print(f"Sending notification to {recipient_role}: {message}")
        # Changed username to "system"
        send_notification(recipient_role, message, "system")

        # Send notification to trainer
        recipient_role = "trainer"
        message = f"New feedback received: {feedback}"
        print(f"Sending notification to {recipient_role}: {message}")
        # Changed username to "system"
        send_notification(recipient_role, message, "system")

        # Send notification to employee
        recipient_role = "employee"
        message = f"New feedback received: {feedback}"
        print(f"Sending notification to {recipient_role}: {message}")
        # Changed username to "system"
        send_notification(recipient_role, message, "system")

        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def generate_pdf_document(questions):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    # Add questions
    for i, question in enumerate(questions, 1):
        pdf.multi_cell(0, 10, txt=f"{i}. {question}", align='L')
        pdf.ln(5)

    # Save to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        pdf.output(temp_file.name)  # Save the PDF to a temporary file
        temp_file.seek(0)  # Move to the beginning of the file
        buffer = temp_file.read()  # Read the content of the file into a buffer

    return buffer


def generate_docx_document(questions):
    doc = docx.Document()

    # Add title
    doc.add_heading('Question Bank', 0)

    # Add questions
    for i, question in enumerate(questions, 1):
        doc.add_paragraph(f"{i}. {question}")
        doc.add_paragraph()  # Add spacing

    # Save to buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_pptx_document(questions):
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Question Bank"

    # Add questions (one per slide)
    for i, question in enumerate(questions, 1):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        title_shape.text = f'Question {i}'

        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = question

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_csv_document(questions):
    # Create CSV in memory
    buffer = io.StringIO()
    for i, question in enumerate(questions, 1):
        buffer.write(f"{i},{question}\n")

    return buffer.getvalue().encode()


def add_questions_to_question_bank(qb_id, topic_name, new_questions):
    db = create_connection()
    if db is None:
        return False

    try:
        # Retrieve existing questions from the question bank
        result = db.question_banks.find_one({"_id": qb_id}, {"questions": 1})

        if result and 'questions' in result:
            existing_questions = result['questions'].split(
                '\n')  # Split existing questions into a list
            updated_questions = existing_questions + new_questions  # Append new questions
            # Join the updated list into a string
            updated_questions_str = '\n'.join(updated_questions)

            # Update the question bank with the new list of questions
            update_result = db.question_banks.update_one(
                {"_id": qb_id},
                {"$set": {"questions": updated_questions_str}}
            )
            return update_result.modified_count > 0
        else:
            return False
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def save_generated_questions(questions, topic_name):
    db = create_connection()
    if db is None:
        return False

    try:
        generated_doc = {
            "topic_name": topic_name,
            "questions": json.dumps(questions),
            "created_at": datetime.now()
        }
        db.generated_questions_history.insert_one(generated_doc)
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def get_generated_questions_history():
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.generated_questions_history.find(
            {}, {"topic_name": 1, "questions": 1, "_id": 0})
        topic_questions = {}
        for result in results_cursor:
            topic_name = result.get('topic_name')
            questions = json.loads(result.get('questions', '[]'))
            if topic_name:
                topic_questions[topic_name] = questions
        return topic_questions
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_questions_from_prompt(prompt, question_type, difficulty_level, num_questions, topic_name):
    try:
        response = model.generate_content(
            f"Generate {num_questions} {question_type} questions based on the following text with {difficulty_level} difficulty:\n\n{prompt}")
        generated_text = response.text
        questions = [line.strip()
                     for line in generated_text.split('\n') if line.strip()]

        # Save generated questions to history
        # Call to save the generated questions
        save_generated_questions(questions, topic_name)

        return questions
    except Exception as e:
        st.error(f"Error generating questions: {e}")
        return None


def get_curriculum_history():
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.curriculum.find({}, {"technology": 1, "topics": 1})
        results = []
        for doc in results_cursor:
            doc['id'] = str(doc['_id'])  # Convert ObjectId to string for 'id'
            results.append(doc)
        return results
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_question_bank_document(qb_id):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.question_banks.find_one(
            {"_id": ObjectId(qb_id)}, {"questions": 1})

        if result and 'questions' in result:
            questions = result['questions'].split('\n')
            doc = docx.Document()
            for question in questions:
                doc.add_paragraph(question)
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_personalized_learning_plan_document(plan_id):
    """
    Generates a DOCX document for a specific learning plan.
    """
    db = create_connection()
    if db is None:
        return None

    try:
        # Fetch the specific learning plan by its unique ID
        learning_plan = db.learning_plans.find_one({"_id": ObjectId(plan_id)})

        if learning_plan:
            doc = docx.Document()
            doc.add_heading(
                f"Learning Plan for {learning_plan.get('username', 'user')}", 0)
            doc.add_paragraph(
                f"Technology: {learning_plan.get('technology', 'N/A')}")
            doc.add_paragraph(
                f"Start Date: {learning_plan.get('start_date', 'N/A')}")
            doc.add_paragraph(
                f"Target End Date: {learning_plan.get('end_date', 'N/A')}")
            doc.add_paragraph(f"Status: {learning_plan.get('status', 'N/A')}")

            # Display time in hours for readability
            estimated_minutes = learning_plan.get('estimated_time', 0)
            estimated_hours = round(estimated_minutes / 60, 1)
            doc.add_paragraph(f"Estimated Time: ~{estimated_hours} hours")

            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            st.error("Learning plan not found.")
            return None
    except Exception as err:
        st.error(f"Database error or invalid ID: {err}")
        return None


def get_all_feedback():
    db = create_connection()
    if db is None:
        return []

    try:
        feedback_records_cursor = db.feedback.find(
            {}, {"username": 1, "question_bank_id": 1, "feedback_text": 1, "rating": 1, "sentiment": 1, "feedback_type": 1, "_id": 0})
        feedback_records = list(feedback_records_cursor)
        # Convert ObjectId to string for question_bank_id if it exists
        for record in feedback_records:
            if 'question_bank_id' in record and isinstance(record['question_bank_id'], ObjectId):
                record['question_bank_id'] = str(record['question_bank_id'])
        return feedback_records
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def generate_feedback_report_document(qb_id):
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.feedback.find({"question_bank_id": ObjectId(qb_id)}, {
                                          "feedback_text": 1, "rating": 1, "_id": 0})
        results = list(results_cursor)

        if results:
            doc = docx.Document()
            for result in results:
                doc.add_paragraph(result.get('feedback_text', ''))
                doc.add_paragraph(str(result.get('rating', 'N/A')))
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_curriculum_mapping_document(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one(
            {"technology": technology}, {"topics": 1, "_id": 0})

        if result and 'topics' in result:
            topics = result['topics']
            # topics is already a string in MongoDB, no need for bytes decoding
            topics = topics.split(',')  # Now split the string into a list
            doc = docx.Document()
            for topic in topics:
                doc.add_paragraph(topic)
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_issue_resolution_report_document(qb_id):
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.issue_resolution.find({"question_bank_id": ObjectId(qb_id)}, {
                                                  "issue": 1, "resolution": 1, "_id": 0})
        results = list(results_cursor)

        if results:
            doc = docx.Document()
            for result in results:
                doc.add_paragraph(result.get('issue', ''))
                doc.add_paragraph(result.get('resolution', ''))
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def generate_assessment_completion_report_document(username):
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.assessments.find(
            {"username": username}, {"score": 1, "completed_at": 1, "_id": 0})
        results = list(results_cursor)

        if results:
            doc = docx.Document()
            for result in results:
                doc.add_paragraph(str(result.get('score', 'N/A')))
                doc.add_paragraph(str(result.get('completed_at', 'N/A')))
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def send_notification(recipient_role, message, username):
    db = create_connection()
    if db is None:
        return False

    try:
        notification_doc = {
            "recipient_role": recipient_role,
            "message": message,
            "username": username,
            "is_new": 1,
            "created_at": datetime.now()
        }
        db.notifications.insert_one(notification_doc)
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def display_notifications(notifications, username):
    if notifications:
        st.sidebar.write("Notifications:")
        for notification in notifications:
            # Check if the notification is new (MongoDB stores 1/0)
            if notification.get('is_new', 0) == 1:
                # Display with a green dot
                st.sidebar.markdown(
                    f"<span style='color: green;'>•</span> {notification.get('message', '')}", unsafe_allow_html=True)
            else:
                st.sidebar.write(notification.get('message', ''))

            # Mark as read when clicked (This interaction might need a separate button/logic in Streamlit)
            # For now, this part is commented out as it requires a specific Streamlit callback.
            # if st.sidebar.button(f"Mark as Read - {notification.get('id', '')}"):
            #     mark_notification_as_read(notification.get('_id')) # Use _id for MongoDB
            #     st.experimental_rerun()  # Refresh the app to show updated notifications

    else:
        st.sidebar.write("No notifications available.")


def mark_notification_as_read(notification_id):
    db = create_connection()
    if db is None:
        return False

    try:
        # Convert notification_id to ObjectId if it's a string
        if isinstance(notification_id, str):
            notification_id = ObjectId(notification_id)

        result = db.notifications.update_one(
            {"_id": notification_id}, {"$set": {"is_new": 0}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def clear_all_notifications(username):
    db = create_connection()
    if db is None:
        return False

    try:
        result = db.notifications.update_many(
            {"username": username}, {"$set": {"is_new": 0}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def get_notifications(recipient_role, username=None):
    db = create_connection()
    if db is None:
        return []

    try:
        query_filter = {"recipient_role": recipient_role}
        if username:
            query_filter["username"] = username

        notifications_cursor = db.notifications.find(
            query_filter).sort("created_at", -1)
        notifications = list(notifications_cursor)
        # Convert ObjectId to string for 'id' if needed for display
        for notification in notifications:
            notification['id'] = str(notification['_id'])
        return notifications
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def clear_notifications(recipient_role, username):
    db = create_connection()
    if db is None:
        return False

    try:
        result = db.notifications.delete_many(
            {"recipient_role": recipient_role, "username": username})
        # Debugging line
        print(
            f"Deleted {result.deleted_count} notifications for {recipient_role} and user {username}.")
        return result.deleted_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# New function to update question bank


def update_question_bank(qb_id, new_questions):
    db = create_connection()
    if db is None:
        return False

    try:
        # Convert qb_id to ObjectId if it's a string
        if isinstance(qb_id, str):
            qb_id = ObjectId(qb_id)

        result = db.question_banks.update_one(
            {"_id": qb_id}, {"$set": {"questions": new_questions}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


# MODIFIED: This version updates status by the plan's unique ID.
# Replace the old 'update_learning_plan_status' function with this one.
def update_learning_plan_status(plan_id, new_status):
    """Updates the status of a specific learning plan document."""
    db = create_connection()
    if db is None:
        return False
    try:
        result = db.learning_plans.update_one(
            {"_id": ObjectId(plan_id)},
            {"$set": {"status": new_status, "updated_at": datetime.now()}}
        )
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


# Replace your entire existing employee_dashboard with this corrected version

def employee_dashboard(username):
    st.header(f"Welcome, {username}! 🚀")
    notifications = get_notifications("employee", username)
    user_id = st.session_state.get('user', {}).get('_id', "unique_user_id")

    # Initialize session state variables for chat and assessment
    if user_id and f"chat_history_{user_id}" not in st.session_state:
        st.session_state[f"chat_history_{user_id}"] = []
    if 'assessment_started' not in st.session_state:
        st.session_state.assessment_started = False
    if 'assessment_finished' not in st.session_state:
        st.session_state.assessment_finished = False

    with st.sidebar:
        selected_tab = option_menu(
            menu_title="Employee Dashboard",
            options=["Your Learning Plan", "Prepare from Material", "Take Assessment", "Your Progress",
                     "Discussion Forum", "Resources", "Chatbot", "Feedback"],
            icons=["book", "file-earmark-text", "check-circle", "bar-chart",
                   "chat", "folder", "chat", "pencil-square"],
            menu_icon="cast", default_index=0, orientation="vertical",
        )

    # --- Your Learning Plan Tab ---
    if selected_tab == "Your Learning Plan":
        st.subheader("Your Learning Plan 📋")
        learning_plans = get_user_learning_plans(username)
        if not learning_plans:
            st.info("You don't have any learning plans assigned yet. Please contact your trainer.")
        else:
            st.info(f"You have {len(learning_plans)} learning plan(s).")
            for plan in learning_plans:
                expander_title = f"{plan['technology']} ({plan['difficulty']}) | Status: {plan['status']}"
                with st.expander(expander_title):
                    st.markdown(f"**Technology:** {plan['technology']}")
                    st.markdown(f"**Difficulty:** {plan['difficulty']}")
                    st.markdown(f"**Assigned On:** {plan['start_date']}")
                    st.markdown(f"**Target End Date:** {plan['end_date']}")
                    estimated_hours = round(plan['estimated_time'] / 60, 1)
                    st.markdown(f"**Estimated Time:** ~{estimated_hours} hours")

                    status_options = ["Assigned", "In Progress", "Completed"]
                    current_status_index = status_options.index(plan.get('status', 'Assigned'))
                    new_status = st.selectbox(
                        "Update Your Progress", options=status_options, index=current_status_index,
                        key=f"status_{plan['_id']}"
                    )
                    if st.button("Update Status", key=f"update_{plan['_id']}"):
                        if update_learning_plan_status(plan['_id'], new_status):
                            st.success("Status updated successfully!")
                            st.rerun()
                        else:
                            st.error("Failed to update status.")
                    
                    st.write("---")
                    plan_doc_data = generate_personalized_learning_plan_document(plan['_id'])
                    if plan_doc_data:
                        st.download_button(
                            label="Download Plan as DOCX",
                            data=plan_doc_data,
                            file_name=f"learning_plan_{plan['technology']}_{plan['_id']}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key=f"download_{plan['_id']}"
                        )

    # --- Prepare from Material Tab ---
    # In employee_dashboard(), update the "Prepare from Material" block

    # In employee_dashboard(), replace the "Prepare from Material" block

    elif selected_tab == "Prepare from Material":
        st.subheader("Prepare from Material 📚")
        # This function correctly fetches only the plans assigned to this user
        learning_plans = get_user_learning_plans(username)
        
        if not learning_plans:
            st.info("No materials are available because no learning plans have been assigned to you.")
        else:
            # Create a list of technologies from the user's assigned plans
            user_technologies = sorted(list(set([plan['technology'] for plan in learning_plans])))
            
            selected_tech = st.selectbox(
                "Select a Curriculum from Your Learning Plan",
                options=user_technologies
            )

            if selected_tech:
                # We need a way to get the curriculum content for the selected technology
                # We will assume a function `get_curriculum_content_by_tech` exists for this
                curriculum_content = get_curriculum_text(selected_tech) 
                
                if curriculum_content:
                    st.markdown("---")
                    st.subheader("Curriculum Content")
                    st.info(curriculum_content)
                    # (Your translation feature can be added here)
                else:
                    st.error("Could not retrieve content for the selected curriculum.")

    # --- Interactive Assessment Tab ---
    # In employee_dashboard(), replace the entire "Take Assessment" block with this:

    # In employee_dashboard(), replace the entire "Take Assessment" block with this:

    # In your employee_dashboard function, replace the entire "Take Assessment" block with this:

    # In your employee_dashboard function, replace the entire "Take Assessment" block with this:

    elif selected_tab == "Take Assessment":
        st.subheader("Take Assessment ✍️")

        # --- View 1: Assessment Selection ---
        if 'assessment_started' not in st.session_state or not st.session_state.assessment_started:
            # Initialize or reset state when starting a new assessment selection
            st.session_state.assessment_started = False
            st.session_state.assessment_finished = False
            for key in ['qb_details', 'questions', 'correct_answers', 'user_answers', 'current_question_index', 'start_time', 'end_time']:
                if key in st.session_state:
                    del st.session_state[key]

            learning_plans = get_user_learning_plans(username)
            available_assessments = [
                p for p in learning_plans if p.get('status') != 'Completed']

            if not available_assessments:
                st.info("You have no available assessments to take.")
            else:
                assessment_options = {
                    p['question_bank_id']: f"{p['technology']} - {p['difficulty']}" for p in available_assessments}
                selected_qb_id_str = st.selectbox(
                    "Select an Assessment from Your Learning Plan",
                    options=list(assessment_options.keys()),
                    format_func=lambda x: assessment_options[x],
                    index=None,
                    placeholder="Choose an assessment..."
                )

                if selected_qb_id_str:
                    qb_id = ObjectId(selected_qb_id_str)
                    all_qbs = get_all_question_banks()
                    qb_details = next(
                        (qb for qb in all_qbs if qb['_id'] == qb_id), None)

                    if qb_details:
                        questions = [q for q in qb_details.get(
                            'questions', '').split('\n') if q.strip()]
                        # 1.5 minutes per question
                        time_limit_minutes = len(questions) * 1.5

                        st.info(f"""
                        **You have selected:** {qb_details['technology']} ({qb_details['difficulty']})
                        - **Number of Questions:** {len(questions)}
                        - **Time Limit:** {int(time_limit_minutes)} minutes
                        """)

                        if st.button("Start Assessment", type="primary"):
                            import random
                            
                            # Get original data
                            original_questions = questions
                            original_correct_answers = get_correct_answers(qb_id)
                            original_options_text = qb_details.get('options', '')
                            original_options_per_question = [opt.strip() for opt in original_options_text.split('|||') if opt.strip()]
                            
                            # Create question indices and shuffle them
                            question_indices = list(range(len(original_questions)))
                            random.shuffle(question_indices)
                            
                            # Reorder questions, answers, and options based on shuffled indices
                            shuffled_questions = [original_questions[i] for i in question_indices]
                            shuffled_correct_answers = [original_correct_answers[i] for i in question_indices]
                            shuffled_options_per_question = [original_options_per_question[i] if i < len(original_options_per_question) else "" for i in question_indices]
                            
                            # For multiple choice questions, also shuffle the options within each question
                            if qb_details.get('question_type', '').lower() == "multiple-choice":
                                final_options_per_question = []
                                final_correct_answers = []
                                
                                for i, (correct_letter, options_str) in enumerate(zip(shuffled_correct_answers, shuffled_options_per_question)):
                                    if options_str:
                                        # Parse options for this question
                                        question_options = [opt.strip() for opt in options_str.split('###') if opt.strip()]
                                        
                                        if len(question_options) >= 2:  # Only shuffle if we have multiple options
                                            # Get the correct answer text before shuffling
                                            correct_index = ord(correct_letter) - ord('A') if correct_letter in ['A', 'B', 'C', 'D'] else 0
                                            correct_text = question_options[correct_index] if correct_index < len(question_options) else question_options[0]
                                            
                                            # Shuffle the options
                                            random.shuffle(question_options)
                                            
                                            # Find the new position of the correct answer
                                            new_correct_index = question_options.index(correct_text)
                                            new_correct_letter = chr(ord('A') + new_correct_index)
                                            
                                            # Store shuffled options and new correct answer
                                            final_options_per_question.append('###'.join(question_options))
                                            final_correct_answers.append(new_correct_letter)
                                        else:
                                            # Not enough options to shuffle, keep as is
                                            final_options_per_question.append(options_str)
                                            final_correct_answers.append(correct_letter)
                                    else:
                                        final_options_per_question.append("")
                                        final_correct_answers.append(correct_letter)
                                
                                # Update the options in qb_details
                                shuffled_qb_details = qb_details.copy()
                                shuffled_qb_details['options'] = '|||'.join(final_options_per_question)
                                shuffled_correct_answers = final_correct_answers
                            else:
                                # For non-MCQ, just use the shuffled data as is
                                shuffled_qb_details = qb_details.copy()
                                shuffled_qb_details['options'] = '|||'.join(shuffled_options_per_question)
                            
                            st.session_state.assessment_started = True
                            st.session_state.assessment_finished = False
                            st.session_state.qb_details = shuffled_qb_details
                            st.session_state.questions = shuffled_questions
                            st.session_state.correct_answers = shuffled_correct_answers
                            st.session_state.user_answers = [None] * len(shuffled_questions)
                            st.session_state.current_question_index = 0
                            st.session_state.start_time = datetime.now()
                            st.session_state.end_time = datetime.now() + timedelta(minutes=time_limit_minutes)
                            st.rerun()

        # --- View 2: Assessment in Progress ---
        elif st.session_state.assessment_started and not st.session_state.assessment_finished:
            time_left = st.session_state.end_time - datetime.now()
            if time_left.total_seconds() < 0:
                st.warning(
                    "Time is up! Submitting your assessment automatically.")
                st.session_state.assessment_finished = True
                st.rerun()

            # Persistent Top Bar
            top_cols = st.columns([3, 1])
            top_cols[0].subheader(
                f"Assessment: {st.session_state.qb_details['technology']}")
            top_cols[1].markdown(
                f"**Time Left: <font color='red'>{str(time_left).split('.')[0]}</font>**", unsafe_allow_html=True)
            st.progress(time_left.total_seconds(
            ) / ((st.session_state.end_time - st.session_state.start_time).total_seconds()))

            # Sidebar Navigation
            with st.sidebar:
                st.write("---")
                st.subheader("Question Palette")
                palette_cols = st.columns(5)
                for i in range(len(st.session_state.questions)):
                    col = palette_cols[i % 5]
                    is_current = (i == st.session_state.current_question_index)
                    is_answered = (
                        st.session_state.user_answers[i] is not None and st.session_state.user_answers[i] != "")
                    btn_type = "primary" if is_current else (
                        "secondary" if is_answered else "secondary")
                    if col.button(f"{i+1}", key=f"nav_{i}", use_container_width=True, type=btn_type):
                        st.session_state.current_question_index = i
                        st.rerun()

            # Main Question Area
            q_idx = st.session_state.current_question_index
            st.markdown(
                f"#### Question {q_idx + 1} of {len(st.session_state.questions)}")
            st.write(st.session_state.questions[q_idx])

            def record_answer():
                widget_key = f"q_widget_{q_idx}"
                if widget_key in st.session_state:
                    st.session_state.user_answers[q_idx] = st.session_state[widget_key]

            question_type = st.session_state.qb_details.get(
                'question_type', '').lower()

            # --- FIXED LOGIC TO DISPLAY CORRECT WIDGET ---
            if question_type == "multiple-choice":
                # Split all options by question separator (|||)
                all_options_text = st.session_state.qb_details.get(
                    'options', '')
                options_per_question = [
                    opt.strip() for opt in all_options_text.split('|||') if opt.strip()]

                # Check if we have options for the current question
                if q_idx < len(options_per_question) and options_per_question[q_idx]:
                    # Split the options for current question by ### separator
                    options = [opt.strip() for opt in options_per_question[q_idx].split(
                        '###') if opt.strip()]

                    if options:  # Make sure we have actual options
                        current_answer = st.session_state.user_answers[q_idx]

                        # Convert current answer letter to index
                        current_index = None
                        if current_answer and current_answer in ['A', 'B', 'C', 'D']:
                            current_index = ord(current_answer) - ord('A')
                            if current_index >= len(options):
                                current_index = None
                        
                        # Custom record_answer function for MCQ to store letter instead of text
                        def record_mcq_answer():
                            widget_key = f"q_widget_{q_idx}"
                            if widget_key in st.session_state:
                                selected_text = st.session_state[widget_key]
                                if selected_text in options:
                                    # Convert selected option text to letter (A, B, C, D)
                                    selected_index = options.index(selected_text)
                                    selected_letter = chr(ord('A') + selected_index)
                                    st.session_state.user_answers[q_idx] = selected_letter

                        st.radio("Select your answer:", options,
                                 key=f"q_widget_{q_idx}", index=current_index, on_change=record_mcq_answer)
                    else:
                        st.warning(
                            f"No valid options found for question {q_idx + 1}.")
                        st.session_state.user_answers[q_idx] = ""
                else:
                    st.warning(
                        f"Options for question {q_idx + 1} are missing or malformed.")
                    st.session_state.user_answers[q_idx] = ""

            elif question_type == "fill-in-the-blank":
                st.text_input(
                    "Your answer:", key=f"q_widget_{q_idx}", value=st.session_state.user_answers[q_idx] or "", on_change=record_answer)

            elif question_type == "subjective":
                st.text_area(
                    "Your answer:", key=f"q_widget_{q_idx}", value=st.session_state.user_answers[q_idx] or "", on_change=record_answer)

            else:
                st.error(
                    f"Unknown question type: '{question_type}'. Please contact your trainer.")

            # --- Bottom Navigation ---
            st.write("---")
            nav_cols = st.columns([1, 1, 2, 1, 1])
            if nav_cols[0].button("Previous", use_container_width=True, disabled=(q_idx == 0)):
                st.session_state.current_question_index -= 1
                st.rerun()
            if nav_cols[4].button("Next", use_container_width=True, disabled=(q_idx >= len(st.session_state.questions) - 1)):
                st.session_state.current_question_index += 1
                st.rerun()
            if nav_cols[2].button("Finish & Submit Assessment", type="primary", use_container_width=True):
                st.session_state.assessment_finished = True
                st.rerun()

        # --- View 3: Results Page ---
        elif st.session_state.assessment_finished:
            st.subheader("Assessment Results")
            score = 0
            for i, user_answer in enumerate(st.session_state.user_answers):
                if user_answer is not None and user_answer.strip().upper() == st.session_state.correct_answers[i].strip().upper():
                    score += 1
            total_questions = len(st.session_state.questions)
            percentage = (score / total_questions) * \
                100 if total_questions > 0 else 0

            st.metric("Final Score",
                      f"{score} / {total_questions}", f"{percentage:.2f}%")
            save_assessment_result(
                username, st.session_state.qb_details['_id'], score)

            with st.expander("Review Your Answers", expanded=True):
                # Get options for displaying in results
                all_options_text = st.session_state.qb_details.get('options', '')
                options_per_question = [opt.strip() for opt in all_options_text.split('|||') if opt.strip()]
                
                for i, q in enumerate(st.session_state.questions):
                    st.write(f"**Question {i+1}: {q}**")
                    user_ans = st.session_state.user_answers[i]
                    correct_ans = st.session_state.correct_answers[i]
                    
                    # For multiple choice, show both letter and text
                    if st.session_state.qb_details.get('question_type', '').lower() == "multiple-choice":
                        if i < len(options_per_question) and options_per_question[i]:
                            options = [opt.strip() for opt in options_per_question[i].split('###') if opt.strip()]
                            
                            user_display = f"{user_ans}) {options[ord(user_ans) - ord('A')]}" if user_ans and user_ans in ['A', 'B', 'C', 'D'] and ord(user_ans) - ord('A') < len(options) else user_ans if user_ans else 'Not Answered'
                            correct_display = f"{correct_ans}) {options[ord(correct_ans) - ord('A')]}" if correct_ans and correct_ans in ['A', 'B', 'C', 'D'] and ord(correct_ans) - ord('A') < len(options) else correct_ans
                        else:
                            user_display = user_ans if user_ans else 'Not Answered'
                            correct_display = correct_ans
                    else:
                        user_display = user_ans if user_ans else 'Not Answered'
                        correct_display = correct_ans
                    
                    if user_ans and user_ans.strip().upper() == correct_ans.strip().upper():
                        st.success(f"✔️ Your answer: {user_display}")
                    else:
                        st.error(f"❌ Your answer: {user_display}")
                        st.info(f"Correct answer: {correct_display}")
                    st.write("---")
            if st.button("Take Another Assessment"):
                # Clean up all session state variables related to the assessment
                keys_to_delete = [k for k in st.session_state.keys() if k.startswith(
                    'assessment_') or k.startswith('q_widget_')]
                for key in keys_to_delete + ['qb_details', 'questions', 'correct_answers', 'user_answers', 'current_question_index', 'start_time', 'end_time']:
                    if key in st.session_state:
                        del st.session_state[key]
                st.rerun()

        
    elif selected_tab == "Your Progress":
        st.subheader("Your Progress 📊")
        completed_assessments = get_completed_assessments(username)

        if completed_assessments:
            # Create DataFrame for better display
            df = pd.DataFrame(completed_assessments)

            # Format completion date
            df['completed_at'] = pd.to_datetime(
                df['completed_at']).dt.strftime('%Y-%m-%d %H:%M')

            # Add filter options
            col1, col2, col3 = st.columns(3)
            with col1:
                tech_filter = st.multiselect(
                    "Filter by Technology",
                    options=sorted(df['technology'].unique()),
                    default=[],
                    help="Select one or more technologies to filter"
                )

            with col2:
                difficulty_filter = st.multiselect(
                    "Filter by Difficulty",
                    options=sorted(df['difficulty'].unique()),
                    default=[],
                    help="Select one or more difficulty levels to filter"
                )

            with col3:
                question_type_filter = st.multiselect(
                    "Filter by Question Type",
                    options=sorted(df['question_type'].unique()),
                    default=[],
                    help="Select one or more question types to filter"
                )

            # Apply filters
            filtered_df = df.copy()
            if tech_filter:
                filtered_df = filtered_df[filtered_df['technology'].isin(
                    tech_filter)]
            if difficulty_filter:
                filtered_df = filtered_df[filtered_df['difficulty'].isin(
                    difficulty_filter)]
            if question_type_filter:
                filtered_df = filtered_df[filtered_df['question_type'].isin(
                    question_type_filter)]

            # Display statistics
            st.subheader("Summary Statistics")
            stats_col1, stats_col2, stats_col3, stats_col4 = st.columns(4)

            with stats_col1:
                st.metric(
                    "Total Assessments",
                    len(filtered_df),
                    help="Total number of completed assessments"
                )

            with stats_col2:
                avg_score = filtered_df['percentage'].mean()
                st.metric(
                    "Average Score",
                    f"{avg_score:.1f}%",
                    help="Average score across all completed assessments"
                )

            with stats_col3:
                best_score = filtered_df['percentage'].max()
                st.metric(
                    "Best Score",
                    f"{best_score:.1f}%",
                    help="Your highest score across all assessments"
                )

            with stats_col4:
                recent_trend = filtered_df.head(3)['percentage'].mean()
                st.metric(
                    "Recent Performance",
                    f"{recent_trend:.1f}%",
                    help="Average score of your last 3 assessments"
                )

            # Display the table with proper formatting
            st.subheader("Detailed Progress")

            # Create a styled dataframe
            display_df = filtered_df[[
                'id', 'technology', 'difficulty', 'question_type', 'score',
                'total_questions', 'percentage', 'completed_at'
            ]].rename(columns={
                'id': 'Assessment ID',
                'technology': 'Technology',
                'difficulty': 'Difficulty',
                'question_type': 'Question Type',
                'score': 'Score',
                'total_questions': 'Total Questions',
                'percentage': 'Percentage (%)',
                'completed_at': 'Completed At'
            })

            # Apply styling based on percentage
            def color_percentage(val):
                try:
                    val_num = float(val)
                    if val_num >= 80:
                        return 'background-color: #90EE90'  # Light green
                    elif val_num >= 60:
                        return 'background-color: #FFD700'  # Gold
                    return 'background-color: #FFB6C1'  # Light red
                except:
                    return ''

            styled_df = display_df.style.applymap(
                color_percentage,
                subset=['Percentage (%)']
            )

            st.dataframe(
                styled_df,
                use_container_width=True,
                hide_index=True
            )

            # Add chart to visualize progress over time
            st.subheader("Progress Over Time")
            fig = px.line(
                display_df.sort_values('Completed At'),
                x='Completed At',
                y='Percentage (%)',
                title='Score Trend',
                markers=True
            )
            st.plotly_chart(fig, use_container_width=True)

            # Add download button for the progress report
            csv = display_df.to_csv(index=False)
            st.download_button(
                label="Download Progress Report",
                data=csv,
                file_name="assessment_progress.csv",
                mime="text/csv",
                help="Download your progress report as a CSV file"
            )

        else:
            st.info("You haven't completed any assessments yet.")

    # In your employee_dashboard function:
    elif selected_tab == "Prepare from Material":
        st.subheader("Prepare from Material 📚")
        learning_plans = get_user_learning_plans(username)
        
        if not learning_plans:
            st.info("No materials available because you have no learning plans assigned.")
        else:
            user_technologies = sorted(list(set([p['technology'] for p in learning_plans])))
            selected_tech = st.selectbox("Select a Curriculum from Your Learning Plan", options=user_technologies)

            if selected_tech:
                curriculum_content = get_curriculum_text(selected_tech)
                if curriculum_content:
                    st.markdown("---")
                    st.subheader("Curriculum Content")
                    st.info(curriculum_content)

                    # --- TRANSLATION FEATURE WITH CHUNKING ---
                    st.markdown("---")
                    st.subheader("Translate Material")
                    languages = ["English", "Hindi", "Tamil", "Telugu", "Spanish", "French", "German", "Chinese", "Japanese", "Korean"]
                    selected_language = st.selectbox("Translate to:", languages)

                    if st.button("Translate", key="translate_material"):
                        with st.spinner(f"Translating to {selected_language}..."):
                            try:
                                # Define the maximum length for each chunk
                                max_chunk_size = 4500  # Well below the 5000 limit
                                
                                # Split the text into chunks
                                text_chunks = [curriculum_content[i:i + max_chunk_size] for i in range(0, len(curriculum_content), max_chunk_size)]
                                
                                translated_chunks = []
                                translator = GoogleTranslator(source='auto', target=selected_language.lower())

                                # Translate each chunk and append to the list
                                for chunk in text_chunks:
                                    translated_chunks.append(translator.translate(chunk))
                                
                                # Join the translated chunks back together
                                full_translated_text = " ".join(translated_chunks)

                                st.subheader(f"Translated Content ({selected_language})")
                                st.success(full_translated_text)
                                
                            except Exception as e:
                                st.error(f"Translation failed. Error: {e}")
                    # --- END OF FEATURE ---
                else:
                    st.error("Could not retrieve content for the selected curriculum.")

    elif selected_tab == "Discussion Forum":
        discussion_forum()

    elif selected_tab == "Resources":
        st.subheader("Resources 🌐")
        st.write(
            "Here you can find various resources to help you in your learning journey.")

        # Search bar for technology
        search_skill = st.text_input("Search for a technology or skill:")

        if search_skill:
            resources_html = generate_resources(search_skill)
            st.markdown(resources_html, unsafe_allow_html=True)
        else:
            print("Please enter Technology")
    elif selected_tab == "Chatbot":
        # Display chatbot interface at the top
        st.subheader("Chat with the AI Trainer 🤖")

        # Initialize message state if it doesn't exist
        if "msg" not in st.session_state:
            st.session_state.msg = ""

        # Create a container for the chat interface
        chat_container = st.container()

        # Define avatars
        user_avatar = "https://static.vecteezy.com/system/resources/previews/009/664/418/non_2x/people-user-team-transparent-free-png.png"
        ai_avatar = "https://thumbs.dreamstime.com/b/chatbot-logo-messenger-ai-robot-icon-vector-illustration-277900892.jpg"

        def clear_text():
            st.session_state.msg = st.session_state.user_input
            st.session_state.user_input = ""

        with chat_container:
            # Display chat messages specific to this user
            for chat in st.session_state[f"chat_history_{user_id}"]:
                if chat['role'] == 'assistant':
                    # Chatbot message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;'>"
                        f"<img src='{ai_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-right: 10px;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>AI:</strong> {chat['content']}</div></div>",
                        unsafe_allow_html=True
                    )
                else:
                    # User message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0; justify-content: flex-end;'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {chat['content']}</div>"
                        f"<img src='{user_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-left: 10px;'>"
                        f"</div>",
                        unsafe_allow_html=True
                    )

            # Input field for user to enter a prompt
            st.text_input("Type your message here...", key="user_input",
                          placeholder="Type a message...", on_change=clear_text)

            if st.session_state.msg:  # Only process if there's a message
                # Append user input to this user's chat history
                st.session_state[f"chat_history_{user_id}"].append(
                    {"role": "user", "content": st.session_state.msg})

                try:
                    # Generate AI response
                    prompt = f"You are an AI assistant for trainers. Respond to the following message: {st.session_state.msg}"
                    response = model.generate_content(prompt)

                    # Handle the response based on the model's response format
                    if hasattr(response, 'parts'):
                        ai_response = ''.join(
                            part.text for part in response.parts)
                    else:
                        ai_response = response.candidates[0].content.parts[0].text

                    # Append AI response to this user's chat history
                    st.session_state[f"chat_history_{user_id}"].append(
                        {"role": "assistant", "content": ai_response})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    ai_response = "I apologize, but I encountered an error. Please try again."
                    st.session_state[f"chat_history_{user_id}"].append(
                        {"role": "assistant", "content": ai_response})

                # Clear the message state
                st.session_state.msg = ""
                # Rerun the app to display the new messages
                # st.experimental_rerun()
                st.rerun()

    elif selected_tab == "Feedback":
        st.subheader("Submit Feedback ✍🏻")

        # Select the type of feedback
        feedback_type = st.radio("Select Feedback Type:",
                                 options=["User  Experience", "Question Bank Feedback", "Assessment Feedback"])

        # Select a question bank for feedback if applicable
        question_banks = get_all_question_banks()
        if feedback_type in ["Question Bank Feedback", "Assessment Feedback"] and question_banks:
            selected_qb = st.selectbox(
                "Select Question Bank for Feedback",
                options=[
                    (str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
                format_func=lambda x: f"ID: {x[0]} - {x[1]}",
                key="feedback_qb_select"
            )
        else:
            selected_qb = None  # No question bank needed for User Experience feedback

        feedback_text = st.text_area("Submit your feedback:", height=100)
        rating = st.slider("Rate your experience (1-5):", 1, 5, 3)

        if st.button("Submit Feedback"):
            if feedback_type == "User  Experience":
                question_bank_id = None  # No question bank for user experience feedback
            else:
                question_bank_id = selected_qb[0] if selected_qb else None

            if submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
                st.success("Feedback submitted successfully!")
            else:
                st.error("Failed to submit feedback.")

    notifications = get_notifications(
        "employee", username)  # Pass the username here
    # Display notifications in the sidebar
    display_notifications(notifications, username)
    if notifications:
        st.sidebar.write("Notifications:")
        for notification in notifications:
            st.sidebar.write(notification['message'])
    else:
        st.sidebar.write("No notifications available.")



# Add this function to your file
@st.cache_data
def translate_text(text, target_language):
    """Translates text to a target language, with caching and error handling."""
    if not text:
        return ""
    try:
        # It's good practice to re-initialize the translator
        translator = googletrans.Translator()
        # The library uses language codes (e.g., 'hi' for Hindi)
        lang_code = googletrans.LANGCODES.get(target_language.lower(), target_language.lower())
        
        # Chunking for long texts
        max_chunk_size = 4500
        text_chunks = [text[i:i + max_chunk_size] for i in range(0, len(text), max_chunk_size)]
        translated_chunks = [translator.translate(chunk, dest=lang_code).text for chunk in text_chunks]
        
        return " ".join(translated_chunks)
    except Exception as e:
        return f"Translation Error: {e}"

# Main function for the discussion forum


def discussion_forum():
    st.title("Discussion Forum 💬")

    # Add a refresh icon beside the title
    refresh_icon = '🔄'
    refresh_button = st.button(refresh_icon, key="refresh_forum_button")

    # Section to post a new message
    new_message = st.text_area("Post a new message", height=100)
    if st.button("Post Message", key="post_message_button"):
        if new_message and st.session_state.user['username']:
            if save_message(new_message, st.session_state.user['username']):
                st.success("Message posted successfully!")
            else:
                st.error("Failed to post message.")
        else:
            st.error("Message and username cannot be empty.")

    # Display all messages in reverse order to show the newest first
    messages = get_messages()
    if messages:
        for message in reversed(messages):  # Reverse the order of messages
            # Create a container for each message
            message_container = st.container()
            with message_container:
                # Display the message with the username
                # Check if the message is from the user
                if message.get('username') == st.session_state.user['username']:
                    # User message on the right
                    st.markdown(
                        f"<div style='display: flex; align-items: center; justify-content: flex-end; margin: 5px 0;'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {message.get('message', '')}</div>"
                        f"</div>",
                        unsafe_allow_html=True
                    )
                else:
                    # Other user's message on the left
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>{message.get('username', 'Unknown')}:</strong> {message.get('message', '')}</div></div>",
                        unsafe_allow_html=True
                    )

                # Fetch and display replies
                replies = get_replies(message['_id'])  # Use MongoDB _id
                if replies:
                    for reply in replies:
                        # Display each reply
                        # Check if the reply is from the user
                        if reply.get('username') == st.session_state.user['username']:
                            # User reply on the right
                            st.markdown(
                                f"<div style='display: flex; align-items: center; justify-content: flex-end; margin: 5px 0;'>"
                                f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                                f"<strong>You:</strong> {reply.get('reply_message', '')}</div>"
                                f"</div>",
                                unsafe_allow_html=True
                            )
                        else:
                            # Other user's reply on the left
                            st.markdown(
                                f"<div style='display: flex; align-items: center; margin: 5px 0;'>"
                                f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                                f"<strong>{reply.get('username', 'Unknown')}:</strong> {reply.get('reply_message', '')}</div></div>",
                                unsafe_allow_html=True
                            )

                else:
                    st.write("  No replies yet.")

            # Right column for reply option
            reply_message = st.text_area(
                "Reply", key=f"reply_area_{message['_id']}")
            if st.button("Post Reply", key=f"reply_button_{message['_id']}"):
                if reply_message:
                    if save_reply(message['_id'], reply_message, st.session_state.user['username']):
                        st.success("Reply posted successfully!")
                        st.experimental_rerun()  # Refresh to show new reply
                    else:
                        st.error("Failed to post reply.")
                else:
                    st.error("Reply cannot be empty.")

    else:
        st.info("No messages available.")

    # Refresh the forum if the refresh button is clicked
    if refresh_button:
        st.session_state.needs_refresh = True

    if 'needs_refresh' in st.session_state and st.session_state.needs_refresh:
        st.session_state.needs_refresh = False
        # Perform any necessary actions here, such as updating the messages or replies
        st.write("Forum refreshed!")


def generate_resources(skill):
    # Use os.getenv for API key
    genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(
        f"Provide a list of resources (articles, videos, etc.) related to the skill '{skill}'. Format the response in an HTML table with the following CSS: table {{ width: 100%; border-collapse: collapse; font-family: sans-serif; }} th, td {{ padding: 10px; text-align: left; border-bottom: 1px solid #ddd; }} th {{ background-color: #3498db; color: white; }}")
    return response.text.strip()


def get_replies(message_id):
    db = create_connection()
    if db is None:
        return []

    try:
        replies_cursor = db.replies.find({"message_id": message_id}, {
                                         "reply_message": 1, "username": 1, "_id": 0})
        replies = list(replies_cursor)
        return replies
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def get_resources():
    db = create_connection()
    if db is None:
        return []
    try:
        resources_cursor = db.resources.find({})
        resources = list(resources_cursor)
        return resources
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def filter_resources(category):
    db = create_connection()
    if db is None:
        return []
    try:
        filtered_resources_cursor = db.resources.find({"category": category})
        filtered_resources = list(filtered_resources_cursor)
        return filtered_resources
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def add_widget(widget):
    db = create_connection()
    if db is None:
        return False
    try:
        db.widgets.insert_one({"widget": widget})
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def get_widgets():
    db = create_connection()
    if db is None:
        return []
    try:
        widgets_cursor = db.widgets.find({})
        widgets = list(widgets_cursor)
        return widgets
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def widget_settings(widget_id):
    db = create_connection()
    if db is None:
        return None
    try:
        # Assuming widget_id is the MongoDB _id
        widget = db.widgets.find_one({"_id": ObjectId(widget_id)})
        return widget
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def change_layout(layout):
    db = create_connection()
    if db is None:
        return False
    try:
        result = db.users.update_one({"username": st.session_state.user['username']}, {
                                     "$set": {"layout": layout}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# Database functions


def save_message(message, username):
    db = create_connection()
    if db is None:
        st.error("Failed to connect to the database.")
        return False

    try:
        message_doc = {
            "username": username,
            "message": message,
            "created_at": datetime.now()
        }
        db.messages.insert_one(message_doc)
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")  # Show the specific database error
        return False


def get_messages():
    db = create_connection()
    if db is None:
        return []
    try:
        messages_cursor = db.messages.find({})
        messages = list(messages_cursor)
        return messages
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def save_reply(message_id, reply_message, username):
    db = create_connection()
    if db is None:
        return False

    try:
        reply_doc = {
            "message_id": message_id,  # This should be the ObjectId of the parent message
            "reply_message": reply_message,
            "username": username,
            "created_at": datetime.now()
        }
        db.replies.insert_one(reply_doc)
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def search_messages(search_query):
    db = create_connection()
    if db is None:
        return []
    try:
        # Use regex for case-insensitive search
        search_results_cursor = db.messages.find(
            {"message": {"$regex": search_query, "$options": "i"}})
        search_results = list(search_results_cursor)
        return search_results
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []


def get_generated_questions():
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.generated_question_files.find_one(
            {}, {"questions": 1, "options": 1, "correct_answers": 1, "_id": 0})
        if result:
            return {
                'questions': result.get('questions', ''),
                'options': result.get('options', ''),
                'correct_answers': result.get('correct_answers', '')
            }
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def get_assessment_results(username):
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.assessments.find({"username": username})
        results = list(results_cursor)
        # Convert ObjectId to string for question_bank_id if it exists
        for result in results:
            if 'question_bank_id' in result and isinstance(result['question_bank_id'], ObjectId):
                result['question_bank_id'] = str(result['question_bank_id'])
        return results
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def get_next_question_bank_id(qb_id):
    db = create_connection()
    if db is None:
        return None

    try:
        # Convert qb_id to ObjectId if it's a string
        if isinstance(qb_id, str):
            qb_id_obj = ObjectId(qb_id)
        else:
            qb_id_obj = qb_id

        # Find the next question bank by _id (which is usually ordered chronologically)
        result = db.question_banks.find(
            {"_id": {"$gt": qb_id_obj}}).sort("_id", 1).limit(1)
        result_doc = next(result, None)

        if result_doc:
            return str(result_doc['_id'])  # Return as string
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


# REFACTORED: To create a learning plan, initiated by a trainer.
# In your stt.py file, replace this function
def create_learning_plan(qb_id_obj, username, trainer_username):
    """
    Creates a learning plan record linking a user, a question bank,
    and the trainer who assigned it.
    """
    db = create_connection()
    if db is None:
        return None
    try:
        # Prevent duplicate assignments
        existing_plan = db.learning_plans.find_one({
            "username": username,
            "question_bank_id": qb_id_obj
        })
        if existing_plan:
            st.warning("This learning plan is already assigned to this user.")
            return None

        qb_details = db.question_banks.find_one({"_id": qb_id_obj})
        if not qb_details:
            st.error("Question Bank details not found.")
            return None

        questions = qb_details.get('questions', '').split('\n')
        difficulty = qb_details.get('difficulty')
        technology = qb_details.get('technology', 'N/A')
        num_questions = len([q for q in questions if q.strip()])
        estimated_time = calculate_estimated_time(num_questions, difficulty)
        start_date = datetime.now()
        num_days = (estimated_time / 60) / 4
        end_date = start_date + timedelta(days=max(1, int(num_days)))

        learning_plan_data = {
            'username': username,
            'question_bank_id': qb_id_obj,
            'trainer_username': trainer_username,  # <-- CRITICAL FIX: Add the trainer's username
            'technology': technology,
            'status': 'Assigned',
            'start_date': start_date.strftime('%Y-%m-%d'),
            'end_date': end_date.strftime('%Y-%m-%d'),
            'estimated_time': estimated_time,
            'created_at': datetime.now()
        }
        result = db.learning_plans.insert_one(learning_plan_data)
        return result.inserted_id
    except OperationFailure as err:
        st.error(f"Database error while creating plan: {err}")
        return None


def get_correct_answers(qb_id):
    db = create_connection()
    if db is None:
        return None

    try:
        # Convert qb_id to ObjectId if it's a string
        if isinstance(qb_id, str):
            qb_id_obj = ObjectId(qb_id)
        else:
            qb_id_obj = qb_id

        result = db.question_answers.find_one({"question_bank_id": qb_id_obj}, {
                                              "answer_data": 1, "_id": 0})

        if result and 'answer_data' in result:
            return result['answer_data'].split('\n')
        else:
            return []
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None


def save_assessment_result(username, qb_id, score):
    db = create_connection()
    if db is None:
        return False

    try:
        # Convert qb_id to ObjectId if it's a string
        if isinstance(qb_id, str):
            qb_id_obj = ObjectId(qb_id)
        else:
            qb_id_obj = qb_id

        assessment_doc = {
            "username": username,
            "question_bank_id": qb_id_obj,  # Store as ObjectId
            "score": score,
            "completed_at": datetime.now()
        }
        db.assessments.insert_one(assessment_doc)

        # Send notification based on score
        if score >= 8:
            feedback_message = f"Great job, {username}! Your score of {score}/10 is excellent!"
        elif score >= 5:
            feedback_message = f"Good effort, {username}! Your score of {score}/10 shows progress."
        else:
            feedback_message = f"Keep practicing, {username}. Your score of {score}/10 indicates more review is needed."

        send_notification("employee", feedback_message,
                          username)  # Pass the username here

        return True

    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def get_previous_learning_plan_end_date(username):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.learning_plans.find_one(
            {"username": username, "status": "Completed"},
            {"end_date": 1, "_id": 0},
            sort=[("end_date", -1)]  # Sort by end_date descending
        )

        if result and 'end_date' in result:
            return datetime.strptime(result['end_date'], '%Y-%m-%d')
        else:
            return datetime.now()
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return datetime.now()  # Return current time on error


def get_completed_date(username):
    db = create_connection()
    if db is None:
        return None

    try:
        # In MongoDB, you don't alter tables to add columns. Just ensure the field is present.
        # The 'updated_at' field will be added when a document is updated/inserted with it.

        result = db.learning_plans.find_one(
            {"username": username, "status": "Completed"},
            {"updated_at": 1, "_id": 0},
            sort=[("updated_at", -1)]  # Sort by updated_at descending
        )

        if result and 'updated_at' in result:
            # updated_at is stored as datetime object
            return result['updated_at']
        else:
            return datetime.now()
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return datetime.now()  # Return current time on error


def calculate_estimated_time(num_questions, difficulty):
    if difficulty == 'Easy':
        estimated_time = num_questions * 10  # 10 minutes per question
    elif difficulty == 'Medium':
        estimated_time = num_questions * 20  # 20 minutes per question
    elif difficulty == 'Hard':
        estimated_time = num_questions * 30  # 30 minutes per question
    return estimated_time


def calculate_estimated_end_date(topics, estimated_time, start_date):
    # Calculate the estimated end date based on the topics length and other analysis
    # For example, assume each topic takes 1 day to complete
    num_days = len(topics)
    estimated_end_date = (datetime.strptime(
        start_date, '%Y-%m-%d') + timedelta(days=num_days)).strftime('%Y-%m-%d')
    return estimated_end_date


# Function to load Lottie animation from a URL


def load_lottie_url(url: str):
    try:
        response = requests.get(url)
        if response.status_code == 200:
            return response.json()
        else:
            st.error(
                f"Failed to load Lottie animation. Status code: {response.status_code}")
            return None
    except Exception as e:
        st.error(f"Error loading Lottie animation: {e}")
        return None


def check_admin_exists():
    """Check if an administrator account already exists"""
    db = create_connection()
    if db is None:
        return False

    try:
        admin_count = db.users.count_documents({"role": "Administrator"})
        return admin_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def main():

    if 'user' not in st.session_state:
        st.session_state.user = None

    if st.session_state.user is None:
        # Show the two-column layout only during login/register
        col1, col2 = st.columns(2)

        # Column 1: Embed Lottie animation
        with col1:
            st.components.v1.html(
                """
                <iframe src="https://lottie.host/embed/1b7b20ac-876d-4a6f-82d5-a1b188f88863/6aZt4s4ExJ.json" 
                        width="100%" height="600" frameborder="0" allowfullscreen></iframe>
                """,
                height=600,
            )

        # Column 2: User authentication
        with col2:
            st.title("Automated Question Builder")
            st.write("")
            st.write("")
            st.write("")

            selected = option_menu(
                menu_title=None,
                options=["Login", "Register"],
                icons=["person", "person-plus"],
                menu_icon="cast",
                default_index=0,
                orientation="horizontal",
            )

            if selected == "Login":
                username = st.text_input("Username 👤", key="login_username",
                                         placeholder="Enter your username", help="Your username")
                password = st.text_input("Password 🔑", type="password", key="login_password",
                                         placeholder="Enter your password", help="Your password")
                if st.button("Login", key="login_button"):
                    user = login_user(username, password)
                    if user:
                        st.session_state.user = user
                        st.success("Logged in successfully!")
                        st.rerun()
                    else:
                        st.error("Invalid username or password")

            elif selected == "Register":
                new_email = st.text_input("Email ✉️ ", key="register_email",
                                          placeholder="Enter your email")
                new_username = st.text_input("Username 👤", key="register_username",
                                             placeholder="Choose a username")
                new_password = st.text_input("Password 🔑", type="password",
                                             key="register_password",
                                             placeholder="Choose a password")

                # Check if admin exists before showing admin role option
                admin_exists = check_admin_exists()
                if admin_exists:
                    role_options = ["Trainer", "Employee"]
                    role = st.selectbox(
                        "Role 👨🏻‍💼", role_options, key="register_role")
                else:
                    role_options = ["Administrator", "Trainer", "Employee"]
                    role = st.selectbox(
                        "Role", role_options, key="register_role")
                    if role == "Administrator":
                        st.warning(
                            "You are registering as the system administrator. This role can only be assigned once.")

                if st.button("Register", key="register_button"):
                    # Double check admin existence before registration
                    if role == "Administrator" and check_admin_exists():
                        st.error(
                            "An administrator account already exists. Please select a different role.")
                    else:
                        if register_user(new_email, new_username, new_password, role):
                            st.success(
                                "Registration successful! Please log in.")
                        else:
                            st.error(
                                "Registration failed. Username may already exist.")

    else:
        # Single column layout for logged-in users

        st.sidebar.write(f"Logged in as: {st.session_state.user['username']}")
        if st.sidebar.button("Logout", key="logout_button"):
            st.session_state.user = None
            st.rerun()

        if st.session_state.user['role'] == 'Administrator':
            admin_dashboard()
        elif st.session_state.user['role'] == 'Trainer':
            trainer_dashboard()
        elif st.session_state.user['role'] == 'Employee':
            employee_dashboard(st.session_state.user['username'])


if __name__ == "__main__":
    main()
